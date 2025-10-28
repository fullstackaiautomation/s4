from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import pandas as pd

def add_text_box(slide, left, top, width, height, text, font_size=10, bold=False, align=PP_ALIGN.LEFT, bg_color=None, text_color=None):
    """Helper function to add a text box to a slide"""
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.text = str(text) if text is not None else ''
    text_frame.word_wrap = True

    # Format the text
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.font.size = Pt(font_size)
    paragraph.font.bold = bold
    paragraph.font.name = 'Calibri'

    if text_color:
        paragraph.font.color.rgb = text_color

    # Add background color if specified
    if bg_color:
        fill = text_box.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

    return text_box

def create_bollard_scorecard(slide, left_offset, bollard_data, is_left=True):
    """Create a single bollard scorecard on the slide"""

    # Card width and spacing
    card_width = Inches(4.5)

    # Header section - Image placeholder for user to add image
    # Create a simple rectangle that users can delete and replace with an image
    img_placeholder = slide.shapes.add_shape(
        1,  # Rectangle
        left_offset,
        Inches(0.3),
        Inches(1.5),
        Inches(1.5)
    )
    img_placeholder.fill.solid()
    img_placeholder.fill.fore_color.rgb = RGBColor(240, 240, 240)
    img_placeholder.line.color.rgb = RGBColor(150, 150, 150)
    img_placeholder.line.width = Pt(1)
    img_placeholder.name = f"ImagePlaceholder_{bollard_data.get('sku', 'SKU')}"

    # Add instruction text
    tf = img_placeholder.text_frame
    tf.text = "[Add Image Here]\n\nDelete this box\nand insert image"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(9)
    tf.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
    tf.vertical_anchor = 1  # Middle vertical alignment

    # Bollard Name and SKU (next to image)
    title_left = left_offset + Inches(1.7)
    add_text_box(slide, title_left, Inches(0.3), Inches(2.7), Inches(0.4),
                 bollard_data.get('name', 'Bollard Name'), font_size=13, bold=True)
    add_text_box(slide, title_left, Inches(0.75), Inches(2.7), Inches(0.25),
                 f"SKU: {bollard_data.get('sku', 'N/A')}", font_size=9, bold=False)

    # Opportunity Score
    add_text_box(slide, title_left, Inches(1.05), Inches(2.7), Inches(0.25),
                 f"Opportunity Score: {bollard_data.get('opportunity_score', 'TBD')}",
                 font_size=10, bold=True, text_color=RGBColor(0, 112, 192))

    # Monthly Profit Impact metrics
    add_text_box(slide, title_left, Inches(1.32), Inches(2.7), Inches(0.18),
                 f"Conservative Monthly Profit Impact: {bollard_data.get('conservative_impact', 'TBD')}",
                 font_size=8, bold=False)
    add_text_box(slide, title_left, Inches(1.50), Inches(2.7), Inches(0.18),
                 f"Moderate Monthly Profit Impact: {bollard_data.get('moderate_impact', 'TBD')}",
                 font_size=8, bold=False)
    add_text_box(slide, title_left, Inches(1.68), Inches(2.7), Inches(0.18),
                 f"Optimistic Monthly Profit Impact: {bollard_data.get('optimistic_impact', 'TBD')}",
                 font_size=8, bold=False)

    # Current position tracker
    current_y = Inches(2.0)
    row_height = Inches(0.22)

    # SALES DATA SECTION - TRANSPOSED (Years as rows, metrics as columns)
    add_text_box(slide, left_offset, current_y, card_width, Inches(0.28),
                 "SALES DATA", font_size=11, bold=True,
                 bg_color=RGBColor(68, 114, 196), text_color=RGBColor(255, 255, 255))
    current_y += Inches(0.32)

    # Sales data headers (metrics across the top)
    metrics = ['Orders', 'QTY Sold', 'Sales Total', 'Profit Total', 'Avg Sale Price', 'Avg CPU']
    year_label_width = Inches(0.6)
    metric_width = Inches(0.63)

    # Header row with metric names (gray background, bold)
    add_text_box(slide, left_offset, current_y, year_label_width, row_height, "Year",
                 font_size=7, bold=True, bg_color=RGBColor(217, 217, 217))
    col_x = left_offset + year_label_width
    for metric in metrics:
        add_text_box(slide, col_x, current_y, metric_width, row_height, metric,
                     font_size=7, bold=True, align=PP_ALIGN.CENTER, bg_color=RGBColor(217, 217, 217))
        col_x += metric_width
    current_y += row_height

    # Data rows for each year
    for year in ['2023', '2024', '2025', 'Total']:
        is_total = (year == 'Total')
        bg = RGBColor(217, 217, 217) if is_total else None

        add_text_box(slide, left_offset, current_y, year_label_width, row_height, year,
                     font_size=7, bold=is_total, bg_color=bg)
        col_x = left_offset + year_label_width
        for metric in metrics:
            value = bollard_data.get('sales', {}).get(year, {}).get(metric, 'TBD')
            add_text_box(slide, col_x, current_y, metric_width, row_height, str(value),
                         font_size=7, align=PP_ALIGN.CENTER, bold=is_total, bg_color=bg)
            col_x += metric_width
        current_y += row_height

    current_y += Inches(0.08)

    # CURRENT PRICING & COMPETITION SECTION - Scorecard Box Style
    add_text_box(slide, left_offset, current_y, card_width, Inches(0.28),
                 "CURRENT PRICING & COMPETITION", font_size=11, bold=True,
                 bg_color=RGBColor(112, 173, 71), text_color=RGBColor(255, 255, 255))
    current_y += Inches(0.32)

    # Single row with all pricing metrics in boxes
    pricing_labels = ['Current CPU', 'ZASP CPU', 'Current MSRP', 'Comp. CPU', 'Profit Impact']
    pricing_keys = ['current_cpu', 'zasp_cpu', 'current_msrp', 'competitor_cpu', 'profit_impact']
    pricing_col_width = card_width / len(pricing_labels)

    col_x = left_offset
    for i, (label, key) in enumerate(zip(pricing_labels, pricing_keys)):
        # Create a box for each metric
        box = slide.shapes.add_shape(
            1,  # Rectangle
            col_x,
            current_y,
            pricing_col_width - Inches(0.02),
            Inches(0.5)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(240, 240, 240)
        box.line.color.rgb = RGBColor(112, 173, 71)
        box.line.width = Pt(2)

        # Add label (top)
        tf = box.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        p1.text = label
        p1.alignment = PP_ALIGN.CENTER
        p1.font.size = Pt(7)
        p1.font.bold = True
        p1.font.name = 'Calibri'
        p1.font.color.rgb = RGBColor(0, 0, 0)

        # Add value (bottom)
        p2 = tf.add_paragraph()
        p2.text = str(bollard_data.get('pricing', {}).get(key, 'TBD'))
        p2.alignment = PP_ALIGN.CENTER
        p2.font.size = Pt(8)
        p2.font.bold = False
        p2.font.name = 'Calibri'
        p2.font.color.rgb = RGBColor(0, 0, 0)
        p2.space_before = Pt(3)

        col_x += pricing_col_width

    current_y += Inches(0.58)

    # RECOMMENDATIONS SECTION - Scorecard Box Style (One Row)
    add_text_box(slide, left_offset, current_y, card_width, Inches(0.28),
                 "RECOMMENDATIONS", font_size=11, bold=True,
                 bg_color=RGBColor(255, 192, 0), text_color=RGBColor(0, 0, 0))
    current_y += Inches(0.32)

    recommendation_labels = ['New MSRP', 'New Margin', 'Est. Monthly Vol.', 'Profit Monthly']
    recommendation_keys = ['new_msrp', 'new_margin', 'est_monthly_volume', 'profit_monthly']
    rec_col_width = card_width / len(recommendation_labels)

    col_x = left_offset
    for label, key in zip(recommendation_labels, recommendation_keys):
        # Create a box for each metric
        box = slide.shapes.add_shape(
            1,  # Rectangle
            col_x,
            current_y,
            rec_col_width - Inches(0.02),
            Inches(0.5)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(255, 250, 230)
        box.line.color.rgb = RGBColor(255, 192, 0)
        box.line.width = Pt(2)

        # Add label (top)
        tf = box.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        p1.text = label
        p1.alignment = PP_ALIGN.CENTER
        p1.font.size = Pt(7)
        p1.font.bold = True
        p1.font.name = 'Calibri'
        p1.font.color.rgb = RGBColor(0, 0, 0)

        # Add value (bottom)
        p2 = tf.add_paragraph()
        p2.text = str(bollard_data.get('recommendations', {}).get(key, 'TBD'))
        p2.alignment = PP_ALIGN.CENTER
        p2.font.size = Pt(8)
        p2.font.bold = False
        p2.font.name = 'Calibri'
        p2.font.color.rgb = RGBColor(0, 0, 0)
        p2.space_before = Pt(3)

        col_x += rec_col_width

    current_y += Inches(0.58)

    # ORDER PROJECTIONS SECTION - Scorecard Box Style
    add_text_box(slide, left_offset, current_y, card_width, Inches(0.28),
                 "ORDER PROJECTIONS", font_size=11, bold=True,
                 bg_color=RGBColor(237, 125, 49), text_color=RGBColor(255, 255, 255))
    current_y += Inches(0.32)

    projection_labels = ['First Order Qty', 'Sell Out Time', 'Short Profit Lift', 'Long Profit Lift']
    projection_keys = ['first_order_qty', 'sellout_timeframe', 'short_profit_lift', 'long_profit_lift']
    proj_col_width = card_width / len(projection_labels)

    col_x = left_offset
    for label, key in zip(projection_labels, projection_keys):
        # Create a box for each metric
        box = slide.shapes.add_shape(
            1,  # Rectangle
            col_x,
            current_y,
            proj_col_width - Inches(0.02),
            Inches(0.5)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(255, 235, 215)
        box.line.color.rgb = RGBColor(237, 125, 49)
        box.line.width = Pt(2)

        # Add label (top)
        tf = box.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        p1.text = label
        p1.alignment = PP_ALIGN.CENTER
        p1.font.size = Pt(7)
        p1.font.bold = True
        p1.font.name = 'Calibri'
        p1.font.color.rgb = RGBColor(0, 0, 0)

        # Add value (bottom)
        p2 = tf.add_paragraph()
        p2.text = str(bollard_data.get('projections', {}).get(key, 'TBD'))
        p2.alignment = PP_ALIGN.CENTER
        p2.font.size = Pt(8)
        p2.font.bold = False
        p2.font.name = 'Calibri'
        p2.font.color.rgb = RGBColor(0, 0, 0)
        p2.space_before = Pt(3)

        col_x += proj_col_width

    current_y += Inches(0.58)

    # NOTES SECTION
    add_text_box(slide, left_offset, current_y, card_width, Inches(0.25),
                 "NOTES", font_size=10, bold=True,
                 bg_color=RGBColor(180, 180, 180))
    current_y += Inches(0.28)

    notes_box = slide.shapes.add_textbox(left_offset, current_y, card_width, Inches(0.5))
    notes_box.fill.solid()
    notes_box.fill.fore_color.rgb = RGBColor(245, 245, 245)
    notes_box.line.color.rgb = RGBColor(150, 150, 150)
    notes_frame = notes_box.text_frame
    notes_frame.word_wrap = True
    notes_frame.text = bollard_data.get('notes', 'Add notes here...')
    notes_frame.paragraphs[0].font.size = Pt(8)

def create_scorecard_slides(prs, bollard_pairs):
    """Create multiple scorecard slides with 2 bollards each"""

    for pair in bollard_pairs:
        # Add a blank slide (no title)
        blank_slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_slide_layout)

        # Left bollard
        create_bollard_scorecard(slide, Inches(0.25), pair[0], is_left=True)

        # Right bollard
        if len(pair) > 1:
            create_bollard_scorecard(slide, Inches(5.25), pair[1], is_left=False)

def create_bollard_data(row):
    """Create a bollard data dictionary from a DataFrame row"""
    return {
        'name': row['Product Name'],
        'sku': row['SKU'],
        'opportunity_score': 'TBD',  # Can be calculated based on sales, profit, growth potential
        'conservative_impact': 'TBD',
        'moderate_impact': 'TBD',
        'optimistic_impact': 'TBD',
        'sales': {
            '2023': {
                'Orders': 'TBD',
                'QTY Sold': 'TBD',
                'Sales Total': 'TBD',
                'Profit Total': 'TBD',
                'Avg Sale Price': 'TBD',
                'Avg CPU': 'TBD'
            },
            '2024': {
                'Orders': 'TBD',
                'QTY Sold': 'TBD',
                'Sales Total': 'TBD',
                'Profit Total': 'TBD',
                'Avg Sale Price': 'TBD',
                'Avg CPU': 'TBD'
            },
            '2025': {
                'Orders': 'TBD',
                'QTY Sold': 'TBD',
                'Sales Total': 'TBD',
                'Profit Total': 'TBD',
                'Avg Sale Price': 'TBD',
                'Avg CPU': 'TBD'
            },
            'Total': {
                'Orders': f"{row['Orders']:.1f}" if pd.notna(row['Orders']) else 'TBD',
                'QTY Sold': f"{row['Order Quantity']:.0f}" if pd.notna(row['Order Quantity']) else 'TBD',
                'Sales Total': f"${row['Sales Total']:,.0f}" if pd.notna(row['Sales Total']) else 'TBD',
                'Profit Total': f"${row['Profit Total']:,.0f}" if pd.notna(row['Profit Total']) else 'TBD',
                'Avg Sale Price': f"${row['Sales Total']/row['Order Quantity']:,.2f}" if pd.notna(row['Sales Total']) and row['Order Quantity'] > 0 else 'TBD',
                'Avg CPU': f"${row['Cost Total']/row['Order Quantity']:,.2f}" if pd.notna(row['Cost Total']) and row['Order Quantity'] > 0 else 'TBD'
            }
        },
        'pricing': {
            'current_cpu': f"${row['Cost Total']/row['Order Quantity']:,.2f}" if pd.notna(row['Cost Total']) and row['Order Quantity'] > 0 else 'TBD',
            'zasp_cpu': 'TBD',
            'current_msrp': 'TBD',
            'competitor_cpu': 'TBD',
            'profit_impact': 'TBD'
        },
        'recommendations': {
            'new_msrp': 'TBD',
            'new_margin': 'TBD',
            'est_monthly_volume': 'TBD',
            'profit_monthly': 'TBD'
        },
        'projections': {
            'first_order_qty': 'TBD',
            'sellout_timeframe': 'TBD',
            'short_profit_lift': 'TBD',
            'long_profit_lift': 'TBD'
        },
        'notes': 'Add product-specific notes and strategic considerations here.'
    }

def main():
    # Load the most recent presentation
    prs_path = r"C:\Users\blkw\OneDrive\Documents\Claude Code\Source 4 Industries\Bulk Bollards\Bulk_Bollard_Order_Proposal_FINAL_v2.pptx"
    prs = Presentation(prs_path)

    # Load bollard data from CSV
    df = pd.read_csv(r"C:\Users\blkw\OneDrive\Documents\Claude Code\Source 4 Industries\Bulk Bollards\S4_Bollards_35_Month_Data.csv")

    # Define the specific bollard pairs by SKU patterns or product name patterns
    bollard_pairs_config = [
        # Pair 1: 4" and 6" Stainless Steel Removable
        [
            {'pattern': '4".*Removable.*Stainless Steel', 'type': 'name'},
            {'pattern': '6".*Removable.*Stainless Steel', 'type': 'name'}
        ],
        # Pair 2: 4" and 6" Carbon Steel Removable
        [
            {'pattern': '4".*Removable.*Carbon Steel', 'type': 'name'},
            {'pattern': '6".*Removable.*Carbon Steel', 'type': 'name'}
        ],
        # Pair 3: 4" Carbon Steel Retractable and 4" Stainless Steel Retractable
        [
            {'pattern': '4".*Retractable.*Carbon Steel', 'type': 'name'},
            {'pattern': '4".*Retractable.*Stainless Steel', 'type': 'name'}
        ],
        # Pair 4: 4" and 6" Baseplate Carbon Steel
        [
            {'pattern': 'BCSV404036', 'type': 'sku'},
            {'pattern': 'BCSV604036', 'type': 'sku'}
        ],
        # Pair 5: 4" and 6" Baseplate Stainless Steel
        [
            {'pattern': 'BSSV404036', 'type': 'sku'},
            {'pattern': 'BSSV604036', 'type': 'sku'}
        ]
    ]

    # Find matching bollards for each pair
    bollard_pairs = []
    for pair_config in bollard_pairs_config:
        pair = []
        for search_config in pair_config:
            pattern = search_config['pattern']
            search_type = search_config['type']

            # Find bollards matching the pattern
            if search_type == 'sku':
                matching = df[df['SKU'] == pattern]
            else:  # search by name
                matching = df[df['Product Name'].str.contains(pattern, case=False, na=False, regex=True)]

            if not matching.empty:
                # Take the one with highest sales
                best_match = matching.nlargest(1, 'Sales Total').iloc[0]
                pair.append(create_bollard_data(best_match))

        if pair:  # Only add if we found at least one bollard
            bollard_pairs.append(pair)

    # Create the scorecard slides
    create_scorecard_slides(prs, bollard_pairs)

    # Save the updated presentation
    output_path = r"C:\Users\blkw\OneDrive\Documents\Claude Code\Source 4 Industries\Bulk Bollards\Bulk_Bollard_Order_Proposal_WITH_SCORECARDS.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Added {sum(len(pair) for pair in bollard_pairs)} bollard scorecards across {len(bollard_pairs)} slides")

    # Print what bollards were included
    for i, pair in enumerate(bollard_pairs):
        print(f"\nSlide {i+1}:")
        for bollard in pair:
            print(f"  - {bollard['name']} ({bollard['sku']})")

if __name__ == "__main__":
    main()
