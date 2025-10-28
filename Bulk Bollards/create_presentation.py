from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import pandas as pd

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define colors
DARK_BLUE = RGBColor(31, 73, 125)
LIGHT_BLUE = RGBColor(68, 114, 196)
GREEN = RGBColor(112, 173, 71)
RED = RGBColor(192, 0, 0)
GRAY = RGBColor(127, 127, 127)

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = DARK_BLUE
    title_para.alignment = PP_ALIGN.CENTER

    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = GRAY
    subtitle_para.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Add blue header bar
    header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)

    return slide

# Slide 1: Title Slide
add_title_slide(prs, "BULK BOLLARD ORDER PROPOSAL", "Direct-to-Manufacturer Partnership with ZASP\nSourceFour Industries")

# Slide 2: Executive Summary
slide = add_content_slide(prs, "Executive Summary")
content = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(8.5), Inches(5.5))
tf = content.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Opportunity Overview"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = LIGHT_BLUE
p.space_after = Pt(12)

bullets = [
    "Transitioning from 1-800 Bollards distributor to direct ZASP manufacturing",
    "$140,087 cost savings on 500-unit initial order (62% reduction)",
    "Margin improvement from 29% to 72% on wholesale bollards",
    "$186,190 additional annual profit potential at historical sales rates",
    "Recommended initial order: 198 units totaling $33,599"
]

for bullet in bullets:
    p = tf.add_paragraph()
    p.text = bullet
    p.font.size = Pt(18)
    p.level = 0
    p.space_after = Pt(8)

# Slide 3: Cost Comparison - ZASP vs 1-800
slide = add_content_slide(prs, "Wholesale Pricing Comparison: ZASP vs 1-800")

# Read data
wholesale = pd.read_csv(r"C:\Users\blkw\OneDrive\Documents\Claude Code\Source 4 Industries\Bulk Bollards\Wholesale_Pricing.csv")

# Add table
left = Inches(0.5)
top = Inches(1.5)
width = Inches(9)
height = Inches(5.5)

rows = 11
cols = 6

table = slide.shapes.add_table(rows, cols, left, top, width, height).table

# Set column widths
table.columns[0].width = Inches(3.5)
table.columns[1].width = Inches(1.2)
table.columns[2].width = Inches(1.2)
table.columns[3].width = Inches(1.2)
table.columns[4].width = Inches(1.2)
table.columns[5].width = Inches(0.9)

# Header row
headers = ["Product", "ZASP CPU", "1800 CPU", "Savings", "S4 Margin", "Qty"]
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = DARK_BLUE
    para = cell.text_frame.paragraphs[0]
    para.font.size = Pt(11)
    para.font.bold = True
    para.font.color.rgb = RGBColor(255, 255, 255)

# Data rows
products = wholesale[wholesale['Product Name'] != 'TOTAL COST'].head(10)
for idx, (_, row) in enumerate(products.iterrows(), start=1):
    table.cell(idx, 0).text = row['Product Name'][:35]
    table.cell(idx, 1).text = f"${row['ZASP CPU']:.0f}"
    table.cell(idx, 2).text = f"${row['1800 CPU']:.0f}"
    table.cell(idx, 3).text = f"${row['Cost Savings']:,.0f}"
    table.cell(idx, 4).text = f"{row['S4 Margin']*100:.0f}%"
    table.cell(idx, 5).text = f"{row['Quantity']:.0f}"

    for col in range(6):
        cell = table.cell(idx, col)
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(10)

# Add summary box
summary_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.7), Inches(9), Inches(0.5))
tf = summary_box.text_frame
p = tf.paragraphs[0]
p.text = "TOTAL SAVINGS: $140,087 (62% cost reduction) | ZASP Margin: 72% vs 1800 Margin: 29%"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = GREEN
p.alignment = PP_ALIGN.CENTER

# Slide 4: 35-Month Sales History
slide = add_content_slide(prs, "Sales Data Analysis - 35 Months Historical")

sales_35 = pd.read_csv(r"C:\Users\blkw\OneDrive\Documents\Claude Code\Source 4 Industries\Bulk Bollards\Pricing_w_Sales_Data_(35_Month).csv")

# Summary boxes
box1 = slide.shapes.add_shape(1, Inches(0.75), Inches(1.5), Inches(2.5), Inches(1.2))
box1.fill.solid()
box1.fill.fore_color.rgb = LIGHT_BLUE
tb = box1.text_frame
p = tb.paragraphs[0]
p.text = "Total Units Sold\n1,875 units"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

box2 = slide.shapes.add_shape(1, Inches(3.5), Inches(1.5), Inches(2.5), Inches(1.2))
box2.fill.solid()
box2.fill.fore_color.rgb = GREEN
tb = box2.text_frame
p = tb.paragraphs[0]
p.text = "Monthly Average\n53.6 units/month"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

box3 = slide.shapes.add_shape(1, Inches(6.25), Inches(1.5), Inches(2.5), Inches(1.2))
box3.fill.solid()
box3.fill.fore_color.rgb = DARK_BLUE
tb = box3.text_frame
p = tb.paragraphs[0]
p.text = "Total Revenue\n$1,105,239"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

# Top sellers table
left = Inches(0.75)
top = Inches(3)
width = Inches(8.5)
height = Inches(3.5)

rows = 6
cols = 4

table = slide.shapes.add_table(rows, cols, left, top, width, height).table

headers = ["Product", "Units Sold", "Revenue", "ZASP Extra Profit"]
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = DARK_BLUE
    para = cell.text_frame.paragraphs[0]
    para.font.size = Pt(12)
    para.font.bold = True
    para.font.color.rgb = RGBColor(255, 255, 255)

top5 = sales_35[sales_35['Product Name'] != 'TOTAL'].nlargest(5, 'Quantity Sold (35 Months)')
for idx, (_, row) in enumerate(top5.iterrows(), start=1):
    table.cell(idx, 0).text = row['Product Name'][:40]
    table.cell(idx, 1).text = f"{row['Quantity Sold (35 Months)']:.0f}"
    table.cell(idx, 2).text = f"${row['S4 Actual Sales Total']:,.0f}"
    table.cell(idx, 3).text = f"${row['ZASP Extra Profit']:,.0f}"

    for col in range(4):
        cell = table.cell(idx, col)
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(11)

# Slide 5: Recent 9-Month Sales (2025)
slide = add_content_slide(prs, "Recent Sales Trends - 2025 (9 Months)")

sales_9 = pd.read_csv(r"C:\Users\blkw\OneDrive\Documents\Claude Code\Source 4 Industries\Bulk Bollards\Pricing_w_Sales_Data_(9_Month).csv")

# Summary boxes
box1 = slide.shapes.add_shape(1, Inches(0.75), Inches(1.5), Inches(2.5), Inches(1.2))
box1.fill.solid()
box1.fill.fore_color.rgb = LIGHT_BLUE
tb = box1.text_frame
p = tb.paragraphs[0]
p.text = "Total Units Sold\n421 units"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

box2 = slide.shapes.add_shape(1, Inches(3.5), Inches(1.5), Inches(2.5), Inches(1.2))
box2.fill.solid()
box2.fill.fore_color.rgb = GREEN
tb = box2.text_frame
p = tb.paragraphs[0]
p.text = "Monthly Average\n46.8 units/month"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

box3 = slide.shapes.add_shape(1, Inches(6.25), Inches(1.5), Inches(2.5), Inches(1.2))
box3.fill.solid()
box3.fill.fore_color.rgb = DARK_BLUE
tb = box3.text_frame
p = tb.paragraphs[0]
p.text = "Total Revenue\n$285,445"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

# Top sellers table
left = Inches(0.75)
top = Inches(3)
width = Inches(8.5)
height = Inches(3.5)

rows = 6
cols = 4

table = slide.shapes.add_table(rows, cols, left, top, width, height).table

headers = ["Product", "Units Sold", "Revenue", "ZASP Extra Profit"]
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = DARK_BLUE
    para = cell.text_frame.paragraphs[0]
    para.font.size = Pt(12)
    para.font.bold = True
    para.font.color.rgb = RGBColor(255, 255, 255)

top5_2025 = sales_9[sales_9['Product Name'] != 'TOTAL'].nlargest(5, 'Quantity Sold (9 Months)')
for idx, (_, row) in enumerate(top5_2025.iterrows(), start=1):
    table.cell(idx, 0).text = row['Product Name'][:40]
    table.cell(idx, 1).text = f"{row['Quantity Sold (9 Months)']:.0f}"
    table.cell(idx, 2).text = f"${row['S4 Actual Sales Total']:,.0f}"
    table.cell(idx, 3).text = f"${row['ZASP Extra Profit']:,.0f}"

    for col in range(4):
        cell = table.cell(idx, col)
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(11)

# Slide 6: Recommended Order Quantities
slide = add_content_slide(prs, "Recommended Initial Order Quantities")

rec = pd.read_csv(r"C:\Users\blkw\OneDrive\Documents\Claude Code\Source 4 Industries\Bulk Bollards\recommended_order.csv")

# Add methodology
method_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.3), Inches(8.5), Inches(0.6))
tf = method_box.text_frame
p = tf.paragraphs[0]
p.text = "Methodology: Based on 35-month sales velocity, weighted with recent 2025 trends. 6-month supply for high movers, 4-month for medium, 3-month for slower items."
p.font.size = Pt(11)
p.font.italic = True

# Table
left = Inches(0.5)
top = Inches(2)
width = Inches(9)
height = Inches(4.2)

rows = min(11, len(rec) + 1)
cols = 6

table = slide.shapes.add_table(rows, cols, left, top, width, height).table

table.columns[0].width = Inches(3.5)
table.columns[1].width = Inches(1.2)
table.columns[2].width = Inches(1)
table.columns[3].width = Inches(1.2)
table.columns[4].width = Inches(1.1)
table.columns[5].width = Inches(1)

headers = ["Product", "Avg/Month", "Supply", "Order Qty", "Unit Cost", "Total"]
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = DARK_BLUE
    para = cell.text_frame.paragraphs[0]
    para.font.size = Pt(11)
    para.font.bold = True
    para.font.color.rgb = RGBColor(255, 255, 255)

for idx, (_, row) in enumerate(rec.head(10).iterrows(), start=1):
    table.cell(idx, 0).text = row['Product'][:32]
    table.cell(idx, 1).text = f"{row['Avg Monthly']:.1f}"
    table.cell(idx, 2).text = f"{row['Months Supply']:.0f}mo"
    table.cell(idx, 3).text = f"{row['Recommended Qty']:.0f}"
    table.cell(idx, 4).text = f"${row['ZASP CPU']:.0f}"
    table.cell(idx, 5).text = f"${row['Total Cost']:,.0f}"

    for col in range(6):
        cell = table.cell(idx, col)
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(10)

# Summary
total_units = rec['Recommended Qty'].sum()
total_cost = rec['Total Cost'].sum()

summary_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.7))
tf = summary_box.text_frame
p = tf.paragraphs[0]
p.text = f"RECOMMENDED INITIAL ORDER: {total_units:.0f} units | Total Investment: ${total_cost:,.2f}"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = GREEN
p.alignment = PP_ALIGN.CENTER

# Slide 7: Monthly Impact & ROI
slide = add_content_slide(prs, "Monthly Impact & Annual Projection")

# Create comparison boxes
# Current situation
box1 = slide.shapes.add_shape(1, Inches(0.75), Inches(1.5), Inches(4), Inches(2.5))
box1.fill.solid()
box1.fill.fore_color.rgb = RED
tb = box1.text_frame
tb.word_wrap = True

p = tb.paragraphs[0]
p.text = "Current (1-800 Bollards)"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

p = tb.add_paragraph()
p.text = "\nMonthly Profit: $7,257"
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(255, 255, 255)

p = tb.add_paragraph()
p.text = "Annual Profit: $87,084"
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(255, 255, 255)

p = tb.add_paragraph()
p.text = "Margin: 23%"
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(255, 255, 255)

# Future with ZASP
box2 = slide.shapes.add_shape(1, Inches(5.25), Inches(1.5), Inches(4), Inches(2.5))
box2.fill.solid()
box2.fill.fore_color.rgb = GREEN
tb = box2.text_frame
tb.word_wrap = True

p = tb.paragraphs[0]
p.text = "With ZASP Direct"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

p = tb.add_paragraph()
p.text = "\nMonthly Profit: $22,773"
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(255, 255, 255)

p = tb.add_paragraph()
p.text = "Annual Profit: $273,270"
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(255, 255, 255)

p = tb.add_paragraph()
p.text = "Margin: 72%"
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(255, 255, 255)

# Impact box
impact_box = slide.shapes.add_shape(1, Inches(0.75), Inches(4.5), Inches(8.5), Inches(1.8))
impact_box.fill.solid()
impact_box.fill.fore_color.rgb = DARK_BLUE
tb = impact_box.text_frame

p = tb.paragraphs[0]
p.text = "ANNUAL IMPACT"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

p = tb.add_paragraph()
p.text = "\n+$186,186 Additional Annual Profit"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = GREEN
p.alignment = PP_ALIGN.CENTER

p = tb.add_paragraph()
p.text = "214% Profit Increase | 49% Margin Improvement"
p.font.size = Pt(18)
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

# ROI note
note_box = slide.shapes.add_textbox(Inches(0.75), Inches(6.5), Inches(8.5), Inches(0.5))
tf = note_box.text_frame
p = tf.paragraphs[0]
p.text = "Based on 35-month historical average of 53.6 units/month. Initial investment of $33,599 pays back in < 2 months."
p.font.size = Pt(12)
p.font.italic = True
p.alignment = PP_ALIGN.CENTER

# Slide 8: Competitor Pricing Analysis
slide = add_content_slide(prs, "Competitive Pricing Landscape")

content = slide.shapes.add_textbox(Inches(0.75), Inches(1.3), Inches(8.5), Inches(0.5))
tf = content.text_frame
p = tf.paragraphs[0]
p.text = "Market research conducted on major bollard suppliers and manufacturers (data from industry pricing research)."
p.font.size = Pt(11)
p.font.italic = True

# Competitor table
left = Inches(0.75)
top = Inches(2)
width = Inches(8.5)
height = Inches(4)

rows = 7
cols = 4

table = slide.shapes.add_table(rows, cols, left, top, width, height).table

headers = ["Product Type", "Market Range", "1-800 Bollards", "ZASP (Our Cost)"]
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = DARK_BLUE
    para = cell.text_frame.paragraphs[0]
    para.font.size = Pt(12)
    para.font.bold = True
    para.font.color.rgb = RGBColor(255, 255, 255)

competitor_data = [
    ["4\" Removable Stainless Steel", "$400-$900", "$475", "$187"],
    ["6\" Removable Stainless Steel", "$600-$1,100", "$671", "$257"],
    ["4\" Removable Carbon Steel", "$250-$500", "$291", "$108"],
    ["6\" Removable Carbon Steel", "$300-$650", "$421", "$125"],
    ["4\" Retractable Stainless", "$800-$1,500", "$1,025", "$275"],
    ["4\" Retractable Carbon", "$400-$800", "$550", "$222"]
]

for idx, data in enumerate(competitor_data, start=1):
    for col, value in enumerate(data):
        cell = table.cell(idx, col)
        cell.text = value
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(11)
        if col == 3:  # Highlight our cost
            cell.fill.solid()
            cell.fill.fore_color.rgb = GREEN
            para.font.bold = True
            para.font.color.rgb = RGBColor(255, 255, 255)

# Summary
summary_box = slide.shapes.add_textbox(Inches(0.75), Inches(6.3), Inches(8.5), Inches(0.8))
tf = summary_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "ZASP pricing is 50-75% below market rates and 60-75% below 1-800 Bollards. This allows us to maintain competitive retail pricing while dramatically improving margins."
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

# Slide 9: Next Steps & Recommendation
slide = add_content_slide(prs, "Next Steps & Recommendation")

content = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(8.5), Inches(5.5))
tf = content.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Recommendation"
p.font.size = Pt(26)
p.font.bold = True
p.font.color.rgb = DARK_BLUE
p.space_after = Pt(18)

p = tf.add_paragraph()
p.text = "Proceed with initial bulk order of 198 units from ZASP for $33,599"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = GREEN
p.space_after = Pt(18)

p = tf.add_paragraph()
p.text = "Key Benefits:"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = LIGHT_BLUE
p.space_after = Pt(12)

benefits = [
    "Immediate cost savings of 62% per unit compared to 1-800",
    "Margin improvement from 23% to 72% on wholesale bollards",
    "Additional $186K annual profit at current sales velocity",
    "Investment payback in less than 2 months",
    "198-unit order provides 3-6 months of inventory based on velocity",
    "Positions S4 competitively against market pricing"
]

for benefit in benefits:
    p = tf.add_paragraph()
    p.text = benefit
    p.font.size = Pt(16)
    p.level = 0
    p.space_after = Pt(10)

p = tf.add_paragraph()
p.text = "\nNext Steps: Finalize ZASP agreement and place initial order to capture margin improvement immediately."
p.font.size = Pt(14)
p.font.italic = True
p.space_after = Pt(10)

# Save presentation
output_path = r"C:\Users\blkw\OneDrive\Documents\Claude Code\Source 4 Industries\Bulk Bollards\Bulk_Bollard_Order_Proposal.pptx"
prs.save(output_path)
print(f"Presentation created successfully: {output_path}")
