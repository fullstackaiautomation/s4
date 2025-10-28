from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Define colors
DARK_BLUE = RGBColor(31, 73, 125)
LIGHT_BLUE = RGBColor(68, 114, 196)
GREEN = RGBColor(112, 173, 71)
RED = RGBColor(192, 0, 0)
ORANGE = RGBColor(237, 125, 49)
GRAY = RGBColor(127, 127, 127)

# Load existing presentation
prs = Presentation(r"C:\Users\blkw\OneDrive\Documents\Claude Code\Source 4 Industries\Bulk Bollards\Bulk_Bollard_Order_Proposal.pptx")

def add_content_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Add blue header bar
    header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)

    return slide

# Delete old competitor slide (slide 8 - index 7) and replace with new one
if len(prs.slides) > 7:
    rId = prs.slides._sldIdLst[7].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[7]

# Insert new competitor slide at position 8
slide = add_content_slide(prs, "Detailed Competitor Pricing Comparison")

# Add subtitle
subtitle_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.2), Inches(8.5), Inches(0.3))
tf = subtitle_box.text_frame
p = tf.paragraphs[0]
p.text = "Research across 8 major competitors - verified pricing where available"
p.font.size = Pt(11)
p.font.italic = True
p.font.color.rgb = GRAY

# Create detailed comparison table
left = Inches(0.4)
top = Inches(1.6)
width = Inches(9.2)
height = Inches(5.2)

rows = 11
cols = 6

table = slide.shapes.add_table(rows, cols, left, top, width, height).table

# Set column widths
table.columns[0].width = Inches(2.8)
table.columns[1].width = Inches(1.0)
table.columns[2].width = Inches(1.3)
table.columns[3].width = Inches(1.3)
table.columns[4].width = Inches(1.4)
table.columns[5].width = Inches(1.4)

# Header row
headers = ["Product", "ZASP", "1-800", "Grainger", "U-Line", "Savings vs 1-800"]
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = DARK_BLUE
    para = cell.text_frame.paragraphs[0]
    para.font.size = Pt(10)
    para.font.bold = True
    para.font.color.rgb = RGBColor(255, 255, 255)
    para.alignment = PP_ALIGN.CENTER

# Competitor data with actual findings
competitor_data = [
    ["4\" SS Removable", "$187", "$475", "N/A", "N/A", "61%"],
    ["6\" SS Removable", "$257", "$671", "N/A", "N/A", "62%"],
    ["4\" CS Removable", "$108", "$291", "$701", "$120-125", "63%"],
    ["6\" CS Removable", "$125", "$421", "N/A", "~$150", "70%"],
    ["4\" CS Retractable", "$222", "$550", "N/A", "N/A", "60%"],
    ["4\" SS Retractable", "$275", "$1,025", "N/A", "N/A", "73%"],
    ["4\" CS Baseplate", "$65", "$84", "N/A", "N/A", "23%"],
    ["6\" CS Baseplate", "$100", "$164", "N/A", "N/A", "39%"],
    ["4\" SS Baseplate", "$170", "$283", "N/A", "N/A", "40%"],
    ["6\" SS Baseplate", "$205", "$552", "N/A", "N/A", "63%"]
]

for idx, data in enumerate(competitor_data, start=1):
    for col, value in enumerate(data):
        cell = table.cell(idx, col)
        cell.text = value
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(9)
        para.alignment = PP_ALIGN.CENTER if col > 0 else PP_ALIGN.LEFT

        # Highlight ZASP column in green
        if col == 1:
            cell.fill.solid()
            cell.fill.fore_color.rgb = GREEN
            para.font.bold = True
            para.font.color.rgb = RGBColor(255, 255, 255)

        # Highlight savings column
        if col == 5:
            savings_pct = int(value.replace('%', ''))
            if savings_pct >= 60:
                cell.fill.solid()
                cell.fill.fore_color.rgb = GREEN
                para.font.bold = True
                para.font.color.rgb = RGBColor(255, 255, 255)
            elif savings_pct >= 40:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BLUE
                para.font.bold = True
                para.font.color.rgb = RGBColor(255, 255, 255)

# Summary box
summary_box = slide.shapes.add_textbox(Inches(0.4), Inches(6.9), Inches(9.2), Inches(0.4))
tf = summary_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Note: 1-800 Bollards pricing verified. Grainger carbon steel: $701 (4\" removable). U-Line 4.5\" removable: $120-125. Ideal Shield, Post Guard, Global Industrial, Bollard Barrier, and Reliance Foundry require quotes."
p.font.size = Pt(9)
p.font.italic = True
p.font.color.rgb = GRAY

# Add sources note
sources_box = slide.shapes.add_textbox(Inches(0.4), Inches(7.2), Inches(9.2), Inches(0.2))
tf = sources_box.text_frame
p = tf.paragraphs[0]
p.text = "Sources: 1800bollards.com, grainger.com, uline.com, idealshield.com, globalindustrial.com, postguard.com, bollardbarrier.com, reliance-foundry.com"
p.font.size = Pt(7)
p.font.color.rgb = GRAY

# Save updated presentation with new name
output_path = r"C:\Users\blkw\OneDrive\Documents\Claude Code\Source 4 Industries\Bulk Bollards\Bulk_Bollard_Order_Proposal_v2.pptx"
prs.save(output_path)
print(f"Presentation updated successfully with competitor pricing grid: {output_path}")
