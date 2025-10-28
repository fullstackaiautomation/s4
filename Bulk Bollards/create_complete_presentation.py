from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import pandas as pd

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Modern color palette
PRIMARY = RGBColor(20, 30, 48)
ACCENT = RGBColor(0, 176, 240)
SUCCESS = RGBColor(76, 175, 80)
WARNING = RGBColor(255, 152, 0)
DANGER = RGBColor(244, 67, 54)
LIGHT_BG = RGBColor(248, 249, 250)
TEXT_DARK = RGBColor(33, 37, 41)
TEXT_LIGHT = RGBColor(255, 255, 255)
GRAY = RGBColor(108, 117, 125)

def add_modern_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = PRIMARY
    bg1.line.fill.background()

    accent_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(-1), Inches(3), Inches(12), Inches(3))
    accent_shape.fill.solid()
    accent_shape.fill.fore_color.rgb = ACCENT
    accent_shape.line.fill.background()
    accent_shape.rotation = -5

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(8.4), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.word_wrap = True
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = TEXT_LIGHT
    title_para.alignment = PP_ALIGN.LEFT

    subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(5), Inches(8.4), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.word_wrap = True
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(20)
    subtitle_para.font.color.rgb = TEXT_LIGHT
    subtitle_para.alignment = PP_ALIGN.LEFT
    return slide

def add_modern_content_slide(prs, title, has_accent=True):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg.line.fill.background()

    if has_accent:
        accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(7.5))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = ACCENT
        accent_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY
    title_para.alignment = PP_ALIGN.LEFT

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.05), Inches(2), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    return slide

def add_stat_card(slide, left, top, width, height, title, value, color):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = color
    card.line.fill.background()
    card.adjustments[0] = 0.05

    value_box = slide.shapes.add_textbox(left, Inches(top.inches + 0.25), width, Inches(height.inches * 0.5))
    tf = value_box.text_frame
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TEXT_LIGHT
    p.alignment = PP_ALIGN.CENTER

    title_box = slide.shapes.add_textbox(left, Inches(top.inches + height.inches * 0.65), width, Inches(height.inches * 0.3))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_LIGHT
    p.alignment = PP_ALIGN.CENTER

# Read all data
wholesale = pd.read_csv(r"C:\Users\blkw\OneDrive\Documents\Claude Code\Source 4 Industries\Bulk Bollards\Wholesale_Pricing.csv")
sales_35 = pd.read_csv(r"C:\Users\blkw\OneDrive\Documents\Claude Code\Source 4 Industries\Bulk Bollards\Pricing_w_Sales_Data_(35_Month).csv")
sales_9 = pd.read_csv(r"C:\Users\blkw\OneDrive\Documents\Claude Code\Source 4 Industries\Bulk Bollards\Pricing_w_Sales_Data_(9_Month).csv")
rec = pd.read_csv(r"C:\Users\blkw\OneDrive\Documents\Claude Code\Source 4 Industries\Bulk Bollards\recommended_order.csv")
competitor = pd.read_csv(r"C:\Users\blkw\OneDrive\Documents\Claude Code\Source 4 Industries\Bulk Bollards\S4 Bollards - First Order Data - PRICING VS COMPETITION.csv")
per_order = pd.read_csv(r"C:\Users\blkw\OneDrive\Documents\Claude Code\Source 4 Industries\Bulk Bollards\S4 Bollards - First Order Data - BOLLARDS PER ORDER.csv")
categories = pd.read_csv(r"C:\Users\blkw\OneDrive\Documents\Claude Code\Source 4 Industries\Bulk Bollards\S4 Bollards - First Order Data - BOLLARDS PER ORDER - CATEGORIES.csv")

print("Creating comprehensive presentation with all data...")

# SLIDE 1: Title
add_modern_title_slide(prs, "BULK BOLLARD\nORDER PROPOSAL", "Direct-to-Manufacturer Partnership with ZASP\nSourceFour Industries | October 2025")

# SLIDE 2: Executive Summary
slide = add_modern_content_slide(prs, "Executive Summary")
add_stat_card(slide, Inches(0.6), Inches(1.5), Inches(2.8), Inches(1.1), "Cost Reduction", "62%", SUCCESS)
add_stat_card(slide, Inches(3.6), Inches(1.5), Inches(2.8), Inches(1.1), "Margin Improvement", "29% → 72%", ACCENT)
add_stat_card(slide, Inches(6.6), Inches(1.5), Inches(2.8), Inches(1.1), "Annual Profit Boost", "$186K", SUCCESS)

content = slide.shapes.add_textbox(Inches(0.6), Inches(2.9), Inches(8.8), Inches(4))
tf = content.text_frame
tf.word_wrap = True
points = [
    "• Transitioning from 1-800 Bollards distributor to direct ZASP manufacturing",
    "• $140,087 total cost savings on 500-unit comparison (62% reduction)",
    "• Recommended initial order: 198 units totaling $33,599",
    "• Investment payback period: Less than 2 months",
    "• Proven demand: 1,875 units sold over 35 months"
]
for point in points:
    p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
    p.text = point
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_DARK
    p.space_after = Pt(14)
    p.line_spacing = 1.3

# SLIDE 3: Pricing Comparison
slide = add_modern_content_slide(prs, "Wholesale Pricing: ZASP vs 1-800 Bollards")
left = Inches(0.4)
top = Inches(1.5)
width = Inches(9.2)
height = Inches(5.3)
rows = 11
cols = 7
table = slide.shapes.add_table(rows, cols, left, top, width, height).table
table.columns[0].width = Inches(2.6)
table.columns[1].width = Inches(1.1)
table.columns[2].width = Inches(1.0)
table.columns[3].width = Inches(1.0)
table.columns[4].width = Inches(1.0)
table.columns[5].width = Inches(1.2)
table.columns[6].width = Inches(1.3)

headers = ["Product", "MSRP", "ZASP CPU", "1-800 CPU", "Savings", "1-800 Margin", "ZASP Margin"]
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = PRIMARY
    para = cell.text_frame.paragraphs[0]
    para.font.size = Pt(10)
    para.font.bold = True
    para.font.color.rgb = TEXT_LIGHT
    para.alignment = PP_ALIGN.CENTER

products = wholesale[wholesale['Product Name'] != 'TOTAL COST'].head(10)
for idx, (_, row) in enumerate(products.iterrows(), start=1):
    # Product name with SKU
    cell = table.cell(idx, 0)
    tf = cell.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = row['Product Name'] + " Bollard"
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p = tf.add_paragraph()
    p.text = row['S4 SKU']
    p.font.size = Pt(7)
    p.font.color.rgb = GRAY

    # MSRP
    table.cell(idx, 1).text = f"${row['S4 MSRP']:.2f}"

    # ZASP CPU with decimals
    table.cell(idx, 2).text = f"${row['ZASP CPU']:.2f}"

    # 1-800 CPU with decimals
    table.cell(idx, 3).text = f"${row['1800 CPU']:.2f}"

    # Savings (just the difference)
    savings = row['1800 CPU'] - row['ZASP CPU']
    table.cell(idx, 4).text = f"${savings:.2f}"

    # 1-800 Margin: (MSRP - 1800 CPU) / MSRP
    margin_1800 = (row['S4 MSRP'] - row['1800 CPU']) / row['S4 MSRP']
    table.cell(idx, 5).text = f"{margin_1800*100:.1f}%"

    # ZASP Margin: (MSRP - ZASP CPU) / MSRP
    margin_zasp = (row['S4 MSRP'] - row['ZASP CPU']) / row['S4 MSRP']
    table.cell(idx, 6).text = f"{margin_zasp*100:.1f}%"

    for col in range(7):
        if col == 0:
            continue  # Skip product name, already formatted
        cell = table.cell(idx, col)
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(9)
        para.font.color.rgb = TEXT_DARK
        para.alignment = PP_ALIGN.CENTER
        if idx % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BG
        if col == 2:  # ZASP CPU column
            cell.fill.solid()
            cell.fill.fore_color.rgb = SUCCESS
            para.font.color.rgb = TEXT_LIGHT
            para.font.bold = True
        if col == 6:  # ZASP Margin column
            para.font.bold = True
            para.font.color.rgb = SUCCESS

summary_box = slide.shapes.add_textbox(Inches(0.4), Inches(7), Inches(9.2), Inches(0.4))
tf = summary_box.text_frame
p = tf.paragraphs[0]
p.text = "Total Cost Savings: $140,087 (62%) | ZASP Margin: 72% vs 1-800 Margin: 29%"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = SUCCESS
p.alignment = PP_ALIGN.CENTER

# SLIDE 4: 35-Month Sales
slide = add_modern_content_slide(prs, "35 Month Sales Performance for Sampled SKUs")

# Calculate total revenue and ZASP extra profit for stat cards
total_revenue = sales_35[sales_35['Product Name'] != 'TOTAL']['S4 Actual Sales Total'].sum()
total_zasp_profit = sales_35[sales_35['Product Name'] != 'TOTAL']['ZASP Extra Profit'].sum()

add_stat_card(slide, Inches(0.5), Inches(1.4), Inches(2.2), Inches(0.95), "Total Units", "1,875", ACCENT)
add_stat_card(slide, Inches(2.9), Inches(1.4), Inches(2.2), Inches(0.95), "Monthly Avg", "53.6", SUCCESS)
add_stat_card(slide, Inches(5.3), Inches(1.4), Inches(2.2), Inches(0.95), "Revenue", f"${total_revenue:,.0f}", PRIMARY)
add_stat_card(slide, Inches(7.7), Inches(1.4), Inches(2.2), Inches(0.95), "ZASP Extra Profit", f"${total_zasp_profit:,.0f}", SUCCESS)

left = Inches(0.3)
top = Inches(2.5)
width = Inches(9.4)
height = Inches(4.7)
rows = 12
cols = 7
table = slide.shapes.add_table(rows, cols, left, top, width, height).table

table.columns[0].width = Inches(2.4)
table.columns[1].width = Inches(1.0)
table.columns[2].width = Inches(1.2)
table.columns[3].width = Inches(1.2)
table.columns[4].width = Inches(1.2)
table.columns[5].width = Inches(1.2)
table.columns[6].width = Inches(1.2)

headers = ["Product", "Units Sold", "Avg Units/Mo", "Revenue", "Monthly Revenue", "ZASP Extra Profit", "ZASP Monthly Extra Profit"]
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = PRIMARY
    para = cell.text_frame.paragraphs[0]
    para.font.size = Pt(9)
    para.font.bold = True
    para.font.color.rgb = TEXT_LIGHT
    para.alignment = PP_ALIGN.CENTER

# Get all 10 products (exclude TOTAL row)
all_products = sales_35[sales_35['Product Name'] != 'TOTAL'].head(10)

# Initialize totals
total_units = 0
total_avg_monthly = 0
total_rev = 0
total_monthly_rev = 0
total_zasp = 0
total_zasp_monthly = 0

for idx, (_, row) in enumerate(all_products.iterrows(), start=1):
    # Product name with "Bollard"
    table.cell(idx, 0).text = row['Product Name'] + " Bollard"

    # Units Sold
    units_sold = row['Quantity Sold (35 Months)']
    table.cell(idx, 1).text = f"{units_sold:.0f}"

    # Average Units Per Month (divided by 35)
    avg_monthly = units_sold / 35
    table.cell(idx, 2).text = f"{avg_monthly:.1f}"

    # Revenue
    revenue = row['S4 Actual Sales Total']
    table.cell(idx, 3).text = f"${revenue:,.0f}"

    # Monthly Revenue (divided by 35)
    monthly_revenue = revenue / 35
    table.cell(idx, 4).text = f"${monthly_revenue:,.0f}"

    # ZASP Extra Profit
    zasp_profit = row['ZASP Extra Profit']
    table.cell(idx, 5).text = f"${zasp_profit:,.0f}"

    # ZASP Monthly Extra Profit (divided by 35)
    zasp_monthly = zasp_profit / 35
    table.cell(idx, 6).text = f"${zasp_monthly:,.0f}"

    # Add to totals
    total_units += units_sold
    total_avg_monthly += avg_monthly
    total_rev += revenue
    total_monthly_rev += monthly_revenue
    total_zasp += zasp_profit
    total_zasp_monthly += zasp_monthly

    for col in range(7):
        cell = table.cell(idx, col)
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(9)
        para.font.color.rgb = TEXT_DARK
        para.alignment = PP_ALIGN.CENTER if col > 0 else PP_ALIGN.LEFT
        if idx % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BG

# Add TOTAL row
table.cell(11, 0).text = "TOTAL"
table.cell(11, 1).text = f"{total_units:.0f}"
table.cell(11, 2).text = f"{total_avg_monthly:.1f}"
table.cell(11, 3).text = f"${total_rev:,.0f}"
table.cell(11, 4).text = f"${total_monthly_rev:,.0f}"
table.cell(11, 5).text = f"${total_zasp:,.0f}"
table.cell(11, 6).text = f"${total_zasp_monthly:,.0f}"

for col in range(7):
    cell = table.cell(11, col)
    para = cell.text_frame.paragraphs[0]
    para.font.size = Pt(10)
    para.font.bold = True
    para.font.color.rgb = TEXT_LIGHT
    para.alignment = PP_ALIGN.CENTER if col > 0 else PP_ALIGN.LEFT
    cell.fill.solid()
    cell.fill.fore_color.rgb = PRIMARY

# SLIDE 5: 2025 Trends
slide = add_modern_content_slide(prs, "2025 Sales for Sampled SKUs")

# Calculate total revenue and ZASP extra profit for stat cards
total_revenue_2025 = sales_9[sales_9['Product Name'] != 'TOTAL']['S4 Actual Sales Total'].sum()
total_zasp_profit_2025 = sales_9[sales_9['Product Name'] != 'TOTAL']['ZASP Extra Profit'].sum()

add_stat_card(slide, Inches(0.5), Inches(1.4), Inches(2.2), Inches(0.95), "Total Units", "421", ACCENT)
add_stat_card(slide, Inches(2.9), Inches(1.4), Inches(2.2), Inches(0.95), "Monthly Avg", "46.8", SUCCESS)
add_stat_card(slide, Inches(5.3), Inches(1.4), Inches(2.2), Inches(0.95), "Revenue", f"${total_revenue_2025:,.0f}", PRIMARY)
add_stat_card(slide, Inches(7.7), Inches(1.4), Inches(2.2), Inches(0.95), "ZASP Extra Profit", f"${total_zasp_profit_2025:,.0f}", SUCCESS)

left = Inches(0.3)
top = Inches(2.5)
width = Inches(9.4)
height = Inches(4.7)
rows = 12
cols = 7
table = slide.shapes.add_table(rows, cols, left, top, width, height).table

table.columns[0].width = Inches(2.4)
table.columns[1].width = Inches(1.0)
table.columns[2].width = Inches(1.2)
table.columns[3].width = Inches(1.2)
table.columns[4].width = Inches(1.2)
table.columns[5].width = Inches(1.2)
table.columns[6].width = Inches(1.2)

headers = ["Product", "Units Sold", "Avg Units/Mo", "Revenue", "Monthly Revenue", "ZASP Extra Profit", "ZASP Monthly Extra Profit"]
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = PRIMARY
    para = cell.text_frame.paragraphs[0]
    para.font.size = Pt(9)
    para.font.bold = True
    para.font.color.rgb = TEXT_LIGHT
    para.alignment = PP_ALIGN.CENTER

# Get all 10 products (exclude TOTAL row)
all_products_2025 = sales_9[sales_9['Product Name'] != 'TOTAL'].head(10)

# Initialize totals
total_units_2025 = 0
total_avg_monthly_2025 = 0
total_rev_2025 = 0
total_monthly_rev_2025 = 0
total_zasp_2025 = 0
total_zasp_monthly_2025 = 0

for idx, (_, row) in enumerate(all_products_2025.iterrows(), start=1):
    # Product name with "Bollard"
    table.cell(idx, 0).text = row['Product Name'] + " Bollard"

    # Units Sold
    units_sold = row['Quantity Sold (9 Months)']
    table.cell(idx, 1).text = f"{units_sold:.0f}"

    # Average Units Per Month (divided by 9)
    avg_monthly = units_sold / 9
    table.cell(idx, 2).text = f"{avg_monthly:.1f}"

    # Revenue
    revenue = row['S4 Actual Sales Total']
    table.cell(idx, 3).text = f"${revenue:,.0f}"

    # Monthly Revenue (divided by 9)
    monthly_revenue = revenue / 9
    table.cell(idx, 4).text = f"${monthly_revenue:,.0f}"

    # ZASP Extra Profit
    zasp_profit = row['ZASP Extra Profit']
    table.cell(idx, 5).text = f"${zasp_profit:,.0f}"

    # ZASP Monthly Extra Profit (divided by 9)
    zasp_monthly = zasp_profit / 9
    table.cell(idx, 6).text = f"${zasp_monthly:,.0f}"

    # Add to totals
    total_units_2025 += units_sold
    total_avg_monthly_2025 += avg_monthly
    total_rev_2025 += revenue
    total_monthly_rev_2025 += monthly_revenue
    total_zasp_2025 += zasp_profit
    total_zasp_monthly_2025 += zasp_monthly

    for col in range(7):
        cell = table.cell(idx, col)
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(9)
        para.font.color.rgb = TEXT_DARK
        para.alignment = PP_ALIGN.CENTER if col > 0 else PP_ALIGN.LEFT
        if idx % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BG

# Add TOTAL row
table.cell(11, 0).text = "TOTAL"
table.cell(11, 1).text = f"{total_units_2025:.0f}"
table.cell(11, 2).text = f"{total_avg_monthly_2025:.1f}"
table.cell(11, 3).text = f"${total_rev_2025:,.0f}"
table.cell(11, 4).text = f"${total_monthly_rev_2025:,.0f}"
table.cell(11, 5).text = f"${total_zasp_2025:,.0f}"
table.cell(11, 6).text = f"${total_zasp_monthly_2025:,.0f}"

for col in range(7):
    cell = table.cell(11, col)
    para = cell.text_frame.paragraphs[0]
    para.font.size = Pt(10)
    para.font.bold = True
    para.font.color.rgb = TEXT_LIGHT
    para.alignment = PP_ALIGN.CENTER if col > 0 else PP_ALIGN.LEFT
    cell.fill.solid()
    cell.fill.fore_color.rgb = PRIMARY

# SLIDE 6: Customer Order Behavior (NEW!)
print("Adding Order Behavior slide...")
slide = add_modern_content_slide(prs, "Order Analytics")

add_stat_card(slide, Inches(0.6), Inches(1.4), Inches(2.15), Inches(0.95), "Avg Order Value", "$3,110", SUCCESS)
add_stat_card(slide, Inches(2.9), Inches(1.4), Inches(2.15), Inches(0.95), "Bollards/Order", "4.8", ACCENT)
add_stat_card(slide, Inches(5.2), Inches(1.4), Inches(2.15), Inches(0.95), "Profit Per Order", "$818", PRIMARY)
add_stat_card(slide, Inches(7.45), Inches(1.4), Inches(1.95), Inches(0.95), "Avg Shipping", "$331", WARNING)

left = Inches(0.45)
top = Inches(2.5)
width = Inches(9.1)
height = Inches(2.2)
rows = 5
cols = 8
table = slide.shapes.add_table(rows, cols, left, top, width, height).table

table.columns[0].width = Inches(1.5)
table.columns[1].width = Inches(1.1)
table.columns[2].width = Inches(1.0)
table.columns[3].width = Inches(1.2)
table.columns[4].width = Inches(1.0)
table.columns[5].width = Inches(1.0)
table.columns[6].width = Inches(1.0)
table.columns[7].width = Inches(1.0)

headers = ["Category", "AOV", "Per Order", "Profit Per Order", "Margin", "Orders", "Refund %", "Discount %"]
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = PRIMARY
    para = cell.text_frame.paragraphs[0]
    para.font.size = Pt(9)
    para.font.bold = True
    para.font.color.rgb = TEXT_LIGHT
    para.alignment = PP_ALIGN.CENTER

for idx, (_, row) in enumerate(categories.iterrows(), start=1):
    # Category
    table.cell(idx, 0).text = row['CATEGORIES'].replace(' Bollards', '')

    # AOV with $
    aov_value = row['AOV'].replace('$', '').replace(',', '')
    table.cell(idx, 1).text = f"${aov_value}"

    # Per Order
    table.cell(idx, 2).text = str(row['Bollards Per Order'])

    # Profit Per Order with $
    profit_value = row['Profit Per Order'].replace('$', '').replace(',', '')
    table.cell(idx, 3).text = f"${profit_value}"

    # Margin (Profit Per Order / AOV)
    aov_num = float(aov_value)
    profit_num = float(profit_value)
    margin = (profit_num / aov_num) * 100 if aov_num > 0 else 0
    table.cell(idx, 4).text = f"{margin:.1f}%"

    # Orders
    table.cell(idx, 5).text = str(row['Orders'])

    # Refund %
    table.cell(idx, 6).text = row['Refund %']

    # Discount %
    table.cell(idx, 7).text = row['Avg Discount %']

    for col in range(8):
        cell = table.cell(idx, col)
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(9)
        para.font.color.rgb = TEXT_DARK
        para.alignment = PP_ALIGN.CENTER if col > 0 else PP_ALIGN.LEFT
        if idx % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BG

subtitle2 = slide.shapes.add_textbox(Inches(0.45), Inches(4.9), Inches(9.1), Inches(0.3))
tf = subtitle2.text_frame
p = tf.paragraphs[0]
p.text = "Sampled SKUs Profit Per Order"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = PRIMARY

left = Inches(0.45)
top = Inches(5.25)
width = Inches(9.1)
height = Inches(2.1)
rows = 11
cols = 6
table = slide.shapes.add_table(rows, cols, left, top, width, height).table

table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(1.2)
table.columns[2].width = Inches(1.4)
table.columns[3].width = Inches(1.6)
table.columns[4].width = Inches(1.3)
table.columns[5].width = Inches(1.6)

headers = ["Product SKU", "Per Order", "AOV", "Profit Per Order", "Margin", "Discount %"]
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = PRIMARY
    para = cell.text_frame.paragraphs[0]
    para.font.size = Pt(8)
    para.font.bold = True
    para.font.color.rgb = TEXT_LIGHT
    para.alignment = PP_ALIGN.CENTER

# Get all 10 products (exclude TOTAL row)
all_10_skus = per_order[per_order['MAIN 10 SKUS'] != 'MAIN 10 TOTAL'].copy()

for idx, (_, row) in enumerate(all_10_skus.iterrows(), start=1):
    # Product SKU
    table.cell(idx, 0).text = row['MAIN 10 SKUS'][:18]

    # Per Order
    table.cell(idx, 1).text = str(row['Bollards Per Order'])

    # AOV with $
    aov_value = row['AOV'].replace('$', '').replace(',', '')
    table.cell(idx, 2).text = f"${aov_value}"

    # Profit Per Order with $
    profit_value = row['Profit Per Order'].replace('$', '').replace(',', '')
    table.cell(idx, 3).text = f"${profit_value}"

    # Margin (Profit Per Order / AOV)
    aov_num = float(aov_value)
    profit_num = float(profit_value)
    margin = (profit_num / aov_num) * 100 if aov_num > 0 else 0
    table.cell(idx, 4).text = f"{margin:.1f}%"

    # Discount %
    table.cell(idx, 5).text = row['Avg Discount %']

    for col in range(6):
        cell = table.cell(idx, col)
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(8)
        para.font.color.rgb = TEXT_DARK
        para.alignment = PP_ALIGN.CENTER if col > 0 else PP_ALIGN.LEFT
        if idx % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BG
        if col == 3:
            para.font.bold = True
            para.font.color.rgb = SUCCESS

# SLIDE 7: Recommended Order
print("Adding Recommended Order slide...")
slide = add_modern_content_slide(prs, "Recommended Initial Order")

method_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.6), Inches(0.4))
tf = method_box.text_frame
p = tf.paragraphs[0]
p.text = "Data-driven: 35-month velocity + 2025 trends | 6-month supply (high velocity), 4-month (medium), 3-month (lower)"
p.font.size = Pt(10)
p.font.italic = True
p.font.color.rgb = GRAY

left = Inches(0.6)
top = Inches(1.8)
width = Inches(8.8)
height = Inches(4.5)
rows = 11
cols = 6
table = slide.shapes.add_table(rows, cols, left, top, width, height).table

table.columns[0].width = Inches(3.3)
table.columns[1].width = Inches(1.1)
table.columns[2].width = Inches(0.9)
table.columns[3].width = Inches(1.1)
table.columns[4].width = Inches(1.2)
table.columns[5].width = Inches(1.2)

headers = ["Product", "Avg/Month", "Supply", "Order Qty", "Unit Cost", "Total"]
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = PRIMARY
    para = cell.text_frame.paragraphs[0]
    para.font.size = Pt(11)
    para.font.bold = True
    para.font.color.rgb = TEXT_LIGHT
    para.alignment = PP_ALIGN.CENTER

for idx, (_, row) in enumerate(rec.head(10).iterrows(), start=1):
    table.cell(idx, 0).text = row['Product'][:32]
    table.cell(idx, 1).text = f"{row['Avg Monthly']:.1f}"
    table.cell(idx, 2).text = f"{row['Months Supply']:.0f}mo"
    table.cell(idx, 3).text = f"{row['Recommended Qty']:.0f}"
    table.cell(idx, 4).text = f"${row['ZASP CPU']:.0f}"
    table.cell(idx, 5).text = f"${row['Total Cost']:,.0f}"
    for col in range(6):
        cell = table.cell(idx, col)
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(10)
        para.font.color.rgb = TEXT_DARK
        para.alignment = PP_ALIGN.CENTER if col > 0 else PP_ALIGN.LEFT
        if idx % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BG

total_units = rec['Recommended Qty'].sum()
total_cost = rec['Total Cost'].sum()

summary_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(6.5), Inches(8.8), Inches(0.8))
summary_card.fill.solid()
summary_card.fill.fore_color.rgb = SUCCESS
summary_card.line.fill.background()
summary_card.adjustments[0] = 0.05

summary_box = slide.shapes.add_textbox(Inches(0.6), Inches(6.65), Inches(8.8), Inches(0.5))
tf = summary_box.text_frame
p = tf.paragraphs[0]
p.text = f"RECOMMENDED ORDER: {total_units:.0f} units | Total Investment: ${total_cost:,.2f} | Payback: <2 months"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = TEXT_LIGHT
p.alignment = PP_ALIGN.CENTER

# SLIDE 8: Financial Impact
print("Adding Financial Impact slide...")
slide = add_modern_content_slide(prs, "Financial Impact Analysis")

# Annual Impact Banner at the TOP
impact_banner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.4), Inches(8.6), Inches(1.6))
impact_banner.fill.solid()
impact_banner.fill.fore_color.rgb = PRIMARY
impact_banner.line.fill.background()
impact_banner.adjustments[0] = 0.05

impact_text = slide.shapes.add_textbox(Inches(0.9), Inches(1.6), Inches(8.2), Inches(1.2))
tf = impact_text.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "ANNUAL IMPACT"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = ACCENT
p.alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = "+$186,186 Additional Profit"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = SUCCESS
p.alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = "214% Profit Increase  |  49% Margin Improvement  |  <2 Month ROI"
p.font.size = Pt(14)
p.font.color.rgb = TEXT_LIGHT
p.alignment = PP_ALIGN.CENTER

# Current State Card (moved down)
current_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(3.3), Inches(4.1), Inches(3.5))
current_card.fill.solid()
current_card.fill.fore_color.rgb = DANGER
current_card.line.fill.background()
current_card.adjustments[0] = 0.05

current_content = slide.shapes.add_textbox(Inches(0.9), Inches(3.5), Inches(3.7), Inches(3))
tf = current_content.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "CURRENT STATE"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = TEXT_LIGHT
p.alignment = PP_ALIGN.CENTER
p.space_after = Pt(20)

metrics = [("1-800 Bollards Distributor", ""), ("Monthly Profit", "$7,257"), ("Annual Profit", "$87,084"), ("Gross Margin", "23%")]
for label, value in metrics:
    p = tf.add_paragraph()
    p.text = f"{label}\n{value}" if value else label
    p.font.size = Pt(16) if value else Pt(14)
    p.font.bold = bool(value)
    p.font.color.rgb = TEXT_LIGHT
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(12)

# With ZASP Direct Card (moved down)
future_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(3.3), Inches(4.1), Inches(3.5))
future_card.fill.solid()
future_card.fill.fore_color.rgb = SUCCESS
future_card.line.fill.background()
future_card.adjustments[0] = 0.05

future_content = slide.shapes.add_textbox(Inches(5.4), Inches(3.5), Inches(3.7), Inches(3))
tf = future_content.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "WITH ZASP DIRECT"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = TEXT_LIGHT
p.alignment = PP_ALIGN.CENTER
p.space_after = Pt(20)

metrics = [("Direct Manufacturing", ""), ("Monthly Profit", "$22,773"), ("Annual Profit", "$273,270"), ("Gross Margin", "72%")]
for label, value in metrics:
    p = tf.add_paragraph()
    p.text = f"{label}\n{value}" if value else label
    p.font.size = Pt(16) if value else Pt(14)
    p.font.bold = bool(value)
    p.font.color.rgb = TEXT_LIGHT
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(12)

# SLIDE 9: Market Opportunity & Search Data (NEW!)
print("Adding Market Opportunity slide...")
slide = add_modern_content_slide(prs, "Market Opportunity & Search Demand")

subtitle = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(8.8), Inches(0.3))
tf = subtitle.text_frame
p = tf.paragraphs[0]
p.text = "Keyword research reveals 74,000 monthly searches with explosive growth opportunities"
p.font.size = Pt(11)
p.font.italic = True
p.font.color.rgb = GRAY

add_stat_card(slide, Inches(0.6), Inches(1.6), Inches(2.15), Inches(0.95), "Total Searches/Mo", "74,000", SUCCESS)
add_stat_card(slide, Inches(2.9), Inches(1.6), Inches(2.15), Inches(0.95), "Market Size", "$3.5B", ACCENT)
add_stat_card(slide, Inches(5.2), Inches(1.6), Inches(2.15), Inches(0.95), "2030 Projection", "$4.5B", PRIMARY)
add_stat_card(slide, Inches(7.45), Inches(1.6), Inches(1.95), Inches(0.95), "CAGR", "6-8%", WARNING)

content = slide.shapes.add_textbox(Inches(0.6), Inches(2.6), Inches(8.8), Inches(2.2))
tf = content.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Top Growth Keywords & Opportunities:"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = PRIMARY
p.space_after = Pt(12)

keywords = [
    "• Fixed Bollards: 590/mo (+3,329% YoY) - LOW competition - HIGHEST PRIORITY",
    "• Crash-Rated Bollards: 390/mo (+243% YoY) - Security focus driving demand",
    "• Parking Lot Bollards: 1,000/mo (+48% YoY) - Commercial growth",
    "• Yellow Bollards: 590/mo (+50% YoY) - Safety/visibility trending",
    "• Security Bollards: 1,300/mo (+22% YoY) - High commercial intent",
    "• Removable Parking Bollards: 170/mo (+40% YoY) - Flexible access control"
]

for kw in keywords:
    p = tf.add_paragraph()
    p.text = kw
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_DARK
    p.space_after = Pt(8)

banner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(5), Inches(8.8), Inches(0.9))
banner.fill.solid()
banner.fill.fore_color.rgb = ACCENT
banner.line.fill.background()
banner.adjustments[0] = 0.05

banner_text = slide.shapes.add_textbox(Inches(0.8), Inches(5.15), Inches(8.4), Inches(0.6))
tf = banner_text.text_frame
p = tf.paragraphs[0]
p.text = "Market Shift: Commercial/Parking Growing (+26-48%) | Residential Declining (-55%)"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = TEXT_LIGHT
p.alignment = PP_ALIGN.CENTER

insights = slide.shapes.add_textbox(Inches(0.6), Inches(6.1), Inches(8.8), Inches(1.2))
tf = insights.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Strategic Implications:"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = PRIMARY
p.space_after = Pt(8)

implications = [
    "• Fixed bollards represent unprecedented SEO opportunity aligning with our 4 baseplate products",
    "• Commercial focus matches our wholesale strategy and bulk order approach",
    "• Strong search demand validates proven sales history and supports scaling plan"
]

for imp in implications:
    p = tf.add_paragraph()
    p.text = imp
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_DARK
    p.space_after = Pt(6)

# SLIDE 10: Competitor Pricing
print("Adding Competitor Pricing slide...")
slide = add_modern_content_slide(prs, "Market Competitive Analysis")

subtitle = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(8.8), Inches(0.3))
tf = subtitle.text_frame
p = tf.paragraphs[0]
p.text = "Real competitor pricing: 1-800 Bollards, Global Industrial, Polector, U-Line"
p.font.size = Pt(11)
p.font.italic = True
p.font.color.rgb = GRAY

left = Inches(0.25)
top = Inches(1.55)
width = Inches(9.5)
height = Inches(5.4)
rows = 11
cols = 9
table = slide.shapes.add_table(rows, cols, left, top, width, height).table

table.columns[0].width = Inches(2.2)
table.columns[1].width = Inches(0.85)
table.columns[2].width = Inches(0.85)
table.columns[3].width = Inches(0.85)
table.columns[4].width = Inches(0.85)
table.columns[5].width = Inches(0.85)
table.columns[6].width = Inches(0.85)
table.columns[7].width = Inches(1.0)
table.columns[8].width = Inches(1.2)

headers = ["Product", "ZASP", "S4 MSRP", "1-800", "Global", "Polector", "U-Line", "S4 New Price", "S4 New Margin"]
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = PRIMARY
    para = cell.text_frame.paragraphs[0]
    para.font.size = Pt(9)
    para.font.bold = True
    para.font.color.rgb = TEXT_LIGHT
    para.alignment = PP_ALIGN.CENTER

for idx, (_, row) in enumerate(competitor.iterrows()):
    if pd.isna(row['Product Name']) or row['Product Name'] == '':
        continue

    # Product name with "Bollard"
    table.cell(idx+1, 0).text = row['Product Name'] + " Bollard"

    # ZASP with $
    table.cell(idx+1, 1).text = "$" + str(row['ZASP CPU']).replace('$', '').strip()

    # S4 MSRP with $
    table.cell(idx+1, 2).text = "$" + str(row['S4 MSRP']).replace('$', '').replace(',', '').strip()

    # Competitor prices with $
    price_1800 = str(row['1800 CPU']).replace('$', '').replace(',', '').strip()
    table.cell(idx+1, 3).text = "$" + price_1800

    price_global = str(row['GLOBAL']).replace('$', '').replace(',', '').strip() if pd.notna(row['GLOBAL']) else '-'
    table.cell(idx+1, 4).text = price_global if price_global == '-' else "$" + price_global

    price_polector = str(row['POLECTOR']).replace('$', '').replace(',', '').strip() if pd.notna(row['POLECTOR']) else '-'
    table.cell(idx+1, 5).text = price_polector if price_polector == '-' else "$" + price_polector

    price_uline = str(row['ULINE']).replace('$', '').replace(',', '').strip() if pd.notna(row['ULINE']) else '-'
    table.cell(idx+1, 6).text = price_uline if price_uline == '-' else "$" + price_uline

    # S4 New Price with $
    table.cell(idx+1, 7).text = "$" + str(row['S4 NEW MSRP']).replace('$', '').replace(',', '').strip()

    # S4 New Margin
    margin_value = str(row['S4 NEW MARGIN']).replace('%', '').strip()
    table.cell(idx+1, 8).text = margin_value + "%"

    # Find lowest competitor price for highlighting
    competitor_prices = []
    for price_str in [price_1800, price_global, price_polector, price_uline]:
        if price_str != '-':
            try:
                competitor_prices.append((float(price_str.replace(',', '')), price_str))
            except:
                pass

    lowest_price = min(competitor_prices, key=lambda x: x[0])[1] if competitor_prices else None

    for col in range(9):
        cell = table.cell(idx+1, col)
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(8)
        para.font.color.rgb = TEXT_DARK
        para.alignment = PP_ALIGN.CENTER if col > 0 else PP_ALIGN.LEFT
        if (idx+1) % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BG
        if col == 1:  # ZASP column
            cell.fill.solid()
            cell.fill.fore_color.rgb = SUCCESS
            para.font.bold = True
            para.font.color.rgb = TEXT_LIGHT
        elif col == 7:  # S4 New Price column
            cell.fill.solid()
            cell.fill.fore_color.rgb = ACCENT
            para.font.bold = True
            para.font.color.rgb = TEXT_LIGHT
        elif col in [3, 4, 5, 6]:  # Competitor columns - highlight lowest
            cell_value = table.cell(idx+1, col).text.replace('$', '').replace(',', '').strip()
            if lowest_price and cell_value == lowest_price:
                para.font.bold = True
                para.font.color.rgb = WARNING

footnote = slide.shapes.add_textbox(Inches(0.35), Inches(7.05), Inches(9.3), Inches(0.35))
tf = footnote.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "ZASP pricing enables 50-80% margin improvement while remaining competitive. Our new pricing maintains market positioning with significantly higher profitability."
p.font.size = Pt(9)
p.font.color.rgb = GRAY

# SLIDE 11: Next Steps
print("Adding Next Steps slide...")
slide = add_modern_content_slide(prs, "Recommendation & Next Steps")

rec_banner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.4), Inches(8.6), Inches(1.2))
rec_banner.fill.solid()
rec_banner.fill.fore_color.rgb = SUCCESS
rec_banner.line.fill.background()
rec_banner.adjustments[0] = 0.05

rec_text = slide.shapes.add_textbox(Inches(0.9), Inches(1.6), Inches(8.2), Inches(0.8))
tf = rec_text.text_frame
p = tf.paragraphs[0]
p.text = "PROCEED WITH INITIAL ORDER"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = TEXT_LIGHT
p.alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = "198 units from ZASP for $33,599"
p.font.size = Pt(26)
p.font.bold = True
p.font.color.rgb = TEXT_LIGHT
p.alignment = PP_ALIGN.CENTER

benefits_title = slide.shapes.add_textbox(Inches(0.7), Inches(2.8), Inches(8.6), Inches(0.4))
tf = benefits_title.text_frame
p = tf.paragraphs[0]
p.text = "KEY BENEFITS"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = PRIMARY

benefits_content = slide.shapes.add_textbox(Inches(0.7), Inches(3.3), Inches(8.6), Inches(3))
tf = benefits_content.text_frame
tf.word_wrap = True

benefits = [
    "✓  62% cost reduction vs 1-800 Bollards distributor",
    "✓  Margin improvement from 23% to 72%",
    "✓  Additional $186K annual profit at current velocity",
    "✓  Investment payback in less than 2 months",
    "✓  3-6 months inventory coverage based on product velocity",
    "✓  Strong market demand: 74,000 monthly searches with explosive growth"
]

for benefit in benefits:
    p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
    p.text = benefit
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_DARK
    p.space_after = Pt(10)

next_banner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(6.5), Inches(8.6), Inches(0.7))
next_banner.fill.solid()
next_banner.fill.fore_color.rgb = ACCENT
next_banner.line.fill.background()
next_banner.adjustments[0] = 0.05

next_text = slide.shapes.add_textbox(Inches(0.9), Inches(6.65), Inches(8.2), Inches(0.4))
tf = next_text.text_frame
p = tf.paragraphs[0]
p.text = "Next: Finalize ZASP agreement → Place initial order → Scale to container-level orders"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = TEXT_LIGHT
p.alignment = PP_ALIGN.CENTER

# Save presentation
output_path = r"C:\Users\blkw\OneDrive\Documents\Claude Code\Source 4 Industries\Bulk Bollards\Bulk_Bollard_Order_Proposal_FINAL_v2.pptx"
prs.save(output_path)
print(f"\n✅ COMPLETE presentation created successfully: {output_path}")
print(f"📊 Total slides: 11")
print(f"📁 Includes: Sales data, competitor pricing, order behavior, keyword research, and recommendations")
