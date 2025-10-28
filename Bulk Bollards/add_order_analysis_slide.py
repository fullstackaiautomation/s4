from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import pandas as pd

# Load existing presentation
prs = Presentation(r"C:\Users\blkw\OneDrive\Documents\Claude Code\Source 4 Industries\Bulk Bollards\Bulk_Bollard_Order_Proposal_FINAL_WITH_COMPETITORS.pptx")

# Modern color palette
PRIMARY = RGBColor(20, 30, 48)
ACCENT = RGBColor(0, 176, 240)
SUCCESS = RGBColor(76, 175, 80)
WARNING = RGBColor(255, 152, 0)
LIGHT_BG = RGBColor(248, 249, 250)
TEXT_DARK = RGBColor(33, 37, 41)
TEXT_LIGHT = RGBColor(255, 255, 255)
GRAY = RGBColor(108, 117, 125)

def add_modern_content_slide(prs, title, has_accent=True):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg.line.fill.background()

    if has_accent:
        accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(7.5))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = ACCENT
        accent_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY
    title_para.alignment = PP_ALIGN.LEFT

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.05), Inches(2), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    return slide

def add_stat_card(slide, left, top, width, height, title, value, color):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = color
    card.line.fill.background()
    card.adjustments[0] = 0.05

    value_box = slide.shapes.add_textbox(left, Inches(top.inches + 0.25), width, Inches(height.inches * 0.5))
    tf = value_box.text_frame
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TEXT_LIGHT
    p.alignment = PP_ALIGN.CENTER

    title_box = slide.shapes.add_textbox(left, Inches(top.inches + height.inches * 0.65), width, Inches(height.inches * 0.3))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_LIGHT
    p.alignment = PP_ALIGN.CENTER

# Read the new data
per_order = pd.read_csv(r"C:\Users\blkw\OneDrive\Documents\Claude Code\Source 4 Industries\Bulk Bollards\S4 Bollards - First Order Data - BOLLARDS PER ORDER.csv")
categories = pd.read_csv(r"C:\Users\blkw\OneDrive\Documents\Claude Code\Source 4 Industries\Bulk Bollards\S4 Bollards - First Order Data - BOLLARDS PER ORDER - CATEGORIES.csv")

# Insert new slide after slide 5 (index 5, which will make it slide 6)
# New Slide: Order Behavior Analysis
slide = add_modern_content_slide(prs, "Customer Order Behavior Analysis")

subtitle = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(8.8), Inches(0.3))
tf = subtitle.text_frame
p = tf.paragraphs[0]
p.text = "Understanding how customers buy bollards - key metrics per product category"
p.font.size = Pt(11)
p.font.italic = True
p.font.color.rgb = GRAY

# Category Performance Cards
add_stat_card(slide, Inches(0.6), Inches(1.6), Inches(2.15), Inches(0.85), "Avg Order Value", "$3,110", SUCCESS)
add_stat_card(slide, Inches(2.9), Inches(1.6), Inches(2.15), Inches(0.85), "Bollards/Order", "4.8", ACCENT)
add_stat_card(slide, Inches(5.2), Inches(1.6), Inches(2.15), Inches(0.85), "Profit/Order", "$818", PRIMARY)
add_stat_card(slide, Inches(7.45), Inches(1.6), Inches(1.95), Inches(0.85), "Avg Shipping", "$331", WARNING)

# Categories table
left = Inches(0.6)
top = Inches(2.6)
width = Inches(8.8)
height = Inches(2.2)

rows = 5
cols = 7

table = slide.shapes.add_table(rows, cols, left, top, width, height).table

table.columns[0].width = Inches(1.8)
table.columns[1].width = Inches(1.15)
table.columns[2].width = Inches(1.15)
table.columns[3].width = Inches(1.15)
table.columns[4].width = Inches(1.15)
table.columns[5].width = Inches(1.15)
table.columns[6].width = Inches(1.15)

headers = ["Category", "AOV", "Per Order", "Profit/Order", "Orders", "Refund %", "Discount %"]
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = PRIMARY
    para = cell.text_frame.paragraphs[0]
    para.font.size = Pt(10)
    para.font.bold = True
    para.font.color.rgb = TEXT_LIGHT
    para.alignment = PP_ALIGN.CENTER

for idx, (_, row) in enumerate(categories.iterrows(), start=1):
    table.cell(idx, 0).text = row['CATEGORIES'].replace(' Bollards', '')
    table.cell(idx, 1).text = row['AOV'].replace('$', '').replace(',', '')
    table.cell(idx, 2).text = str(row['Bollards Per Order'])
    table.cell(idx, 3).text = row['Profit Per Order'].replace('$', '').replace(',', '')
    table.cell(idx, 4).text = str(row['Orders'])
    table.cell(idx, 5).text = row['Refund %']
    table.cell(idx, 6).text = row['Avg Discount %']

    for col in range(7):
        cell = table.cell(idx, col)
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(10)
        para.font.color.rgb = TEXT_DARK
        para.alignment = PP_ALIGN.CENTER if col > 0 else PP_ALIGN.LEFT

        if idx % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BG

# Main 10 SKUs performance table
subtitle2 = slide.shapes.add_textbox(Inches(0.6), Inches(5), Inches(8.8), Inches(0.3))
tf = subtitle2.text_frame
p = tf.paragraphs[0]
p.text = "Top 10 Wholesale SKUs - Order Performance Metrics"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = PRIMARY

left = Inches(0.6)
top = Inches(5.35)
width = Inches(8.8)
height = Inches(1.9)

rows = 6
cols = 7

table = slide.shapes.add_table(rows, cols, left, top, width, height).table

table.columns[0].width = Inches(1.8)
table.columns[1].width = Inches(0.9)
table.columns[2].width = Inches(1.2)
table.columns[3].width = Inches(1.2)
table.columns[4].width = Inches(1.2)
table.columns[5].width = Inches(1.2)
table.columns[6].width = Inches(1.2)

headers = ["Product SKU", "Per Order", "AOV", "Sales AOV", "Profit/Order", "Shipping", "Discount %"]
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = PRIMARY
    para = cell.text_frame.paragraphs[0]
    para.font.size = Pt(9)
    para.font.bold = True
    para.font.color.rgb = TEXT_LIGHT
    para.alignment = PP_ALIGN.CENTER

# Get top 5 by profit per order
top5_orders = per_order[per_order['MAIN 10 SKUS'] != 'MAIN 10 TOTAL'].copy()
top5_orders['Profit_Numeric'] = top5_orders['Profit Per Order'].str.replace('$', '').str.replace(',', '').astype(float)
top5_orders = top5_orders.nlargest(5, 'Profit_Numeric')

for idx, (_, row) in enumerate(top5_orders.iterrows(), start=1):
    table.cell(idx, 0).text = row['MAIN 10 SKUS'][:18]
    table.cell(idx, 1).text = str(row['Bollards Per Order'])
    table.cell(idx, 2).text = row['AOV'].replace('$', '').replace(',', '')
    table.cell(idx, 3).text = row['Sales AOV'].replace('$', '').replace(',', '')
    table.cell(idx, 4).text = row['Profit Per Order'].replace('$', '').replace(',', '')
    table.cell(idx, 5).text = row['Avg Shipping'].replace('$', '').replace(',', '')
    table.cell(idx, 6).text = row['Avg Discount %']

    for col in range(7):
        cell = table.cell(idx, col)
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(9)
        para.font.color.rgb = TEXT_DARK
        para.alignment = PP_ALIGN.CENTER if col > 0 else PP_ALIGN.LEFT

        if idx % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BG

        if col == 4:  # Highlight profit per order
            para.font.bold = True
            para.font.color.rgb = SUCCESS

# Move this slide to position 6 (after current slide 5)
# Get the slide we just added (it's the last one)
xml_slides = prs.slides._sldIdLst
new_slide_element = xml_slides[-1]
# Remove it from the end
xml_slides.remove(new_slide_element)
# Insert it at position 5 (0-indexed, so it becomes slide 6)
xml_slides.insert(5, new_slide_element)

# Save presentation
output_path = r"C:\Users\blkw\OneDrive\Documents\Claude Code\Source 4 Industries\Bulk Bollards\Bulk_Bollard_Order_Proposal_FINAL_WITH_COMPETITORS.pptx"
prs.save(output_path)
print(f"Presentation updated with Order Behavior Analysis slide: {output_path}")
