from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import pandas as pd

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Modern color palette
PRIMARY = RGBColor(20, 30, 48)  # Dark navy
ACCENT = RGBColor(0, 176, 240)  # Modern cyan
SUCCESS = RGBColor(76, 175, 80)  # Material green
WARNING = RGBColor(255, 152, 0)  # Material orange
DANGER = RGBColor(244, 67, 54)  # Material red
LIGHT_BG = RGBColor(248, 249, 250)  # Light gray
TEXT_DARK = RGBColor(33, 37, 41)
TEXT_LIGHT = RGBColor(255, 255, 255)
GRAY = RGBColor(108, 117, 125)

def add_modern_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Modern gradient background effect using overlapping shapes
    bg1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = PRIMARY
    bg1.line.fill.background()

    # Accent shape - modern diagonal element
    accent_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(-1), Inches(3), Inches(12), Inches(3))
    accent_shape.fill.solid()
    accent_shape.fill.fore_color.rgb = ACCENT
    accent_shape.line.fill.background()
    accent_shape.rotation = -5

    # Title with modern positioning
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(8.4), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.word_wrap = True
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = TEXT_LIGHT
    title_para.alignment = PP_ALIGN.LEFT

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(5), Inches(8.4), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.word_wrap = True
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(20)
    subtitle_para.font.color.rgb = TEXT_LIGHT
    subtitle_para.alignment = PP_ALIGN.LEFT

    return slide

def add_modern_content_slide(prs, title, has_accent=True):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Clean white background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg.line.fill.background()

    if has_accent:
        # Modern accent bar - thin and sleek
        accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(7.5))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = ACCENT
        accent_bar.line.fill.background()

    # Title with modern styling
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY
    title_para.alignment = PP_ALIGN.LEFT

    # Subtle underline
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.05), Inches(2), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    return slide

def add_stat_card(slide, left, top, width, height, title, value, color):
    """Add a modern card-style statistic"""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = color
    card.line.fill.background()
    card.adjustments[0] = 0.05

    value_box = slide.shapes.add_textbox(left, Inches(top.inches + 0.3), width, Inches(height.inches * 0.5))
    tf = value_box.text_frame
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = TEXT_LIGHT
    p.alignment = PP_ALIGN.CENTER

    title_box = slide.shapes.add_textbox(left, Inches(top.inches + height.inches * 0.65), width, Inches(height.inches * 0.3))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_LIGHT
    p.alignment = PP_ALIGN.CENTER

# Read data
wholesale = pd.read_csv(r"C:\Users\blkw\OneDrive\Documents\Claude Code\Source 4 Industries\Bulk Bollards\Wholesale_Pricing.csv")
sales_35 = pd.read_csv(r"C:\Users\blkw\OneDrive\Documents\Claude Code\Source 4 Industries\Bulk Bollards\Pricing_w_Sales_Data_(35_Month).csv")
sales_9 = pd.read_csv(r"C:\Users\blkw\OneDrive\Documents\Claude Code\Source 4 Industries\Bulk Bollards\Pricing_w_Sales_Data_(9_Month).csv")
rec = pd.read_csv(r"C:\Users\blkw\OneDrive\Documents\Claude Code\Source 4 Industries\Bulk Bollards\recommended_order.csv")
competitor = pd.read_csv(r"C:\Users\blkw\OneDrive\Documents\Claude Code\Source 4 Industries\Bulk Bollards\S4 Bollards - First Order Data - PRICING VS COMPETITION.csv")

# Slide 1: Modern Title Slide
add_modern_title_slide(prs, "BULK BOLLARD\nORDER PROPOSAL", "Direct-to-Manufacturer Partnership with ZASP\nSourceFour Industries | October 2025")

# Slide 2: Executive Summary with Cards
slide = add_modern_content_slide(prs, "Executive Summary")

add_stat_card(slide, Inches(0.6), Inches(1.5), Inches(2.8), Inches(1.1), "Cost Reduction", "62%", SUCCESS)
add_stat_card(slide, Inches(3.6), Inches(1.5), Inches(2.8), Inches(1.1), "Margin Improvement", "29% → 72%", ACCENT)
add_stat_card(slide, Inches(6.6), Inches(1.5), Inches(2.8), Inches(1.1), "Annual Profit Boost", "$186K", SUCCESS)

content = slide.shapes.add_textbox(Inches(0.6), Inches(2.9), Inches(8.8), Inches(4))
tf = content.text_frame
tf.word_wrap = True

points = [
    "• Transitioning from 1-800 Bollards distributor to direct ZASP manufacturing",
    "• $140,087 total cost savings on 500-unit comparison (62% reduction)",
    "• Recommended initial order: 198 units totaling $33,599",
    "• Investment payback period: Less than 2 months",
    "• Proven demand: 1,875 units sold over 35 months"
]

for point in points:
    p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
    p.text = point
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_DARK
    p.space_after = Pt(14)
    p.line_spacing = 1.3

# Slide 3: Modern Pricing Comparison
slide = add_modern_content_slide(prs, "Wholesale Pricing: ZASP vs 1-800 Bollards")

left = Inches(0.6)
top = Inches(1.5)
width = Inches(8.8)
height = Inches(5.3)

rows = 11
cols = 6

table = slide.shapes.add_table(rows, cols, left, top, width, height).table

table.columns[0].width = Inches(3.2)
table.columns[1].width = Inches(1.1)
table.columns[2].width = Inches(1.1)
table.columns[3].width = Inches(1.3)
table.columns[4].width = Inches(1.1)
table.columns[5].width = Inches(1.0)

headers = ["Product", "ZASP", "1-800", "Savings", "Margin", "Qty"]
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = PRIMARY
    para = cell.text_frame.paragraphs[0]
    para.font.size = Pt(11)
    para.font.bold = True
    para.font.color.rgb = TEXT_LIGHT
    para.alignment = PP_ALIGN.CENTER

products = wholesale[wholesale['Product Name'] != 'TOTAL COST'].head(10)
for idx, (_, row) in enumerate(products.iterrows(), start=1):
    table.cell(idx, 0).text = row['Product Name'][:35]
    table.cell(idx, 1).text = f"${row['ZASP CPU']:.0f}"
    table.cell(idx, 2).text = f"${row['1800 CPU']:.0f}"
    table.cell(idx, 3).text = f"${row['Cost Savings']:,.0f}"
    table.cell(idx, 4).text = f"{row['S4 Margin']*100:.0f}%"
    table.cell(idx, 5).text = f"{row['Quantity']:.0f}"

    for col in range(6):
        cell = table.cell(idx, col)
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(10)
        para.font.color.rgb = TEXT_DARK
        para.alignment = PP_ALIGN.CENTER if col > 0 else PP_ALIGN.LEFT

        if idx % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BG

        if col == 1:
            cell.fill.solid()
            cell.fill.fore_color.rgb = SUCCESS
            para.font.color.rgb = TEXT_LIGHT
            para.font.bold = True

summary_box = slide.shapes.add_textbox(Inches(0.6), Inches(7), Inches(8.8), Inches(0.4))
tf = summary_box.text_frame
p = tf.paragraphs[0]
p.text = "Total Cost Savings: $140,087 (62%) | ZASP Margin: 72% vs 1-800 Margin: 29%"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = SUCCESS
p.alignment = PP_ALIGN.CENTER

# Slide 4: 35-Month Sales
slide = add_modern_content_slide(prs, "Sales Performance: 35-Month Historical Data")

add_stat_card(slide, Inches(0.7), Inches(1.4), Inches(2.6), Inches(0.95), "Total Units", "1,875", ACCENT)
add_stat_card(slide, Inches(3.6), Inches(1.4), Inches(2.6), Inches(0.95), "Monthly Avg", "53.6", SUCCESS)
add_stat_card(slide, Inches(6.5), Inches(1.4), Inches(2.6), Inches(0.95), "Revenue", "$1.1M", PRIMARY)

left = Inches(0.7)
top = Inches(2.6)
width = Inches(8.6)
height = Inches(4.2)

rows = 6
cols = 4

table = slide.shapes.add_table(rows, cols, left, top, width, height).table

headers = ["Product", "Units Sold", "Revenue", "ZASP Extra Profit"]
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = PRIMARY
    para = cell.text_frame.paragraphs[0]
    para.font.size = Pt(12)
    para.font.bold = True
    para.font.color.rgb = TEXT_LIGHT

top5 = sales_35[sales_35['Product Name'] != 'TOTAL'].nlargest(5, 'Quantity Sold (35 Months)')
for idx, (_, row) in enumerate(top5.iterrows(), start=1):
    table.cell(idx, 0).text = row['Product Name'][:40]
    table.cell(idx, 1).text = f"{row['Quantity Sold (35 Months)']:.0f}"
    table.cell(idx, 2).text = f"${row['S4 Actual Sales Total']:,.0f}"
    table.cell(idx, 3).text = f"${row['ZASP Extra Profit']:,.0f}"

    for col in range(4):
        cell = table.cell(idx, col)
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(11)
        para.font.color.rgb = TEXT_DARK
        para.alignment = PP_ALIGN.CENTER if col > 0 else PP_ALIGN.LEFT

        if idx % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BG

        if col == 3:
            para.font.bold = True
            para.font.color.rgb = SUCCESS

# Slide 5: 2025 Trends
slide = add_modern_content_slide(prs, "Recent Sales Trends: 2025 (9 Months)")

add_stat_card(slide, Inches(0.7), Inches(1.4), Inches(2.6), Inches(0.95), "Total Units", "421", ACCENT)
add_stat_card(slide, Inches(3.6), Inches(1.4), Inches(2.6), Inches(0.95), "Monthly Avg", "46.8", SUCCESS)
add_stat_card(slide, Inches(6.5), Inches(1.4), Inches(2.6), Inches(0.95), "Revenue", "$285K", PRIMARY)

left = Inches(0.7)
top = Inches(2.6)
width = Inches(8.6)
height = Inches(4.2)

rows = 6
cols = 4

table = slide.shapes.add_table(rows, cols, left, top, width, height).table

headers = ["Product", "Units Sold", "Revenue", "ZASP Extra Profit"]
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = PRIMARY
    para = cell.text_frame.paragraphs[0]
    para.font.size = Pt(12)
    para.font.bold = True
    para.font.color.rgb = TEXT_LIGHT

top5_2025 = sales_9[sales_9['Product Name'] != 'TOTAL'].nlargest(5, 'Quantity Sold (9 Months)')
for idx, (_, row) in enumerate(top5_2025.iterrows(), start=1):
    table.cell(idx, 0).text = row['Product Name'][:40]
    table.cell(idx, 1).text = f"{row['Quantity Sold (9 Months)']:.0f}"
    table.cell(idx, 2).text = f"${row['S4 Actual Sales Total']:,.0f}"
    table.cell(idx, 3).text = f"${row['ZASP Extra Profit']:,.0f}"

    for col in range(4):
        cell = table.cell(idx, col)
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(11)
        para.font.color.rgb = TEXT_DARK
        para.alignment = PP_ALIGN.CENTER if col > 0 else PP_ALIGN.LEFT

        if idx % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BG

        if col == 3:
            para.font.bold = True
            para.font.color.rgb = SUCCESS

# Slide 6: Recommended Order
slide = add_modern_content_slide(prs, "Recommended Initial Order")

method_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.6), Inches(0.4))
tf = method_box.text_frame
p = tf.paragraphs[0]
p.text = "Data-driven methodology: 35-month sales velocity + 2025 trends | 6-month supply for high velocity, 4-month for medium, 3-month for lower velocity"
p.font.size = Pt(10)
p.font.italic = True
p.font.color.rgb = GRAY

left = Inches(0.6)
top = Inches(1.8)
width = Inches(8.8)
height = Inches(4.5)

rows = 11
cols = 6

table = slide.shapes.add_table(rows, cols, left, top, width, height).table

table.columns[0].width = Inches(3.3)
table.columns[1].width = Inches(1.1)
table.columns[2].width = Inches(0.9)
table.columns[3].width = Inches(1.1)
table.columns[4].width = Inches(1.2)
table.columns[5].width = Inches(1.2)

headers = ["Product", "Avg/Month", "Supply", "Order Qty", "Unit Cost", "Total"]
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = PRIMARY
    para = cell.text_frame.paragraphs[0]
    para.font.size = Pt(11)
    para.font.bold = True
    para.font.color.rgb = TEXT_LIGHT
    para.alignment = PP_ALIGN.CENTER

for idx, (_, row) in enumerate(rec.head(10).iterrows(), start=1):
    table.cell(idx, 0).text = row['Product'][:32]
    table.cell(idx, 1).text = f"{row['Avg Monthly']:.1f}"
    table.cell(idx, 2).text = f"{row['Months Supply']:.0f}mo"
    table.cell(idx, 3).text = f"{row['Recommended Qty']:.0f}"
    table.cell(idx, 4).text = f"${row['ZASP CPU']:.0f}"
    table.cell(idx, 5).text = f"${row['Total Cost']:,.0f}"

    for col in range(6):
        cell = table.cell(idx, col)
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(10)
        para.font.color.rgb = TEXT_DARK
        para.alignment = PP_ALIGN.CENTER if col > 0 else PP_ALIGN.LEFT

        if idx % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BG

total_units = rec['Recommended Qty'].sum()
total_cost = rec['Total Cost'].sum()

summary_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(6.5), Inches(8.8), Inches(0.8))
summary_card.fill.solid()
summary_card.fill.fore_color.rgb = SUCCESS
summary_card.line.fill.background()
summary_card.adjustments[0] = 0.05

summary_box = slide.shapes.add_textbox(Inches(0.6), Inches(6.65), Inches(8.8), Inches(0.5))
tf = summary_box.text_frame
p = tf.paragraphs[0]
p.text = f"RECOMMENDED ORDER: {total_units:.0f} units | Total Investment: ${total_cost:,.2f} | Payback: <2 months"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = TEXT_LIGHT
p.alignment = PP_ALIGN.CENTER

# Slide 7: Impact Analysis
slide = add_modern_content_slide(prs, "Financial Impact Analysis")

current_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.5), Inches(4.1), Inches(3.5))
current_card.fill.solid()
current_card.fill.fore_color.rgb = DANGER
current_card.line.fill.background()
current_card.adjustments[0] = 0.05

current_content = slide.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(3.7), Inches(3))
tf = current_content.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "CURRENT STATE"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = TEXT_LIGHT
p.alignment = PP_ALIGN.CENTER
p.space_after = Pt(20)

metrics = [
    ("1-800 Bollards Distributor", ""),
    ("Monthly Profit", "$7,257"),
    ("Annual Profit", "$87,084"),
    ("Gross Margin", "23%")
]

for label, value in metrics:
    p = tf.add_paragraph()
    p.text = f"{label}\n{value}" if value else label
    p.font.size = Pt(16) if value else Pt(14)
    p.font.bold = bool(value)
    p.font.color.rgb = TEXT_LIGHT
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(12)

future_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.5), Inches(4.1), Inches(3.5))
future_card.fill.solid()
future_card.fill.fore_color.rgb = SUCCESS
future_card.line.fill.background()
future_card.adjustments[0] = 0.05

future_content = slide.shapes.add_textbox(Inches(5.4), Inches(1.7), Inches(3.7), Inches(3))
tf = future_content.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "WITH ZASP DIRECT"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = TEXT_LIGHT
p.alignment = PP_ALIGN.CENTER
p.space_after = Pt(20)

metrics = [
    ("Direct Manufacturing", ""),
    ("Monthly Profit", "$22,773"),
    ("Annual Profit", "$273,270"),
    ("Gross Margin", "72%")
]

for label, value in metrics:
    p = tf.add_paragraph()
    p.text = f"{label}\n{value}" if value else label
    p.font.size = Pt(16) if value else Pt(14)
    p.font.bold = bool(value)
    p.font.color.rgb = TEXT_LIGHT
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(12)

impact_banner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(5.3), Inches(8.6), Inches(1.6))
impact_banner.fill.solid()
impact_banner.fill.fore_color.rgb = PRIMARY
impact_banner.line.fill.background()
impact_banner.adjustments[0] = 0.05

impact_text = slide.shapes.add_textbox(Inches(0.9), Inches(5.5), Inches(8.2), Inches(1.2))
tf = impact_text.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "ANNUAL IMPACT"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = ACCENT
p.alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = "+$186,186 Additional Profit"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = SUCCESS
p.alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = "214% Profit Increase  |  49% Margin Improvement  |  <2 Month ROI"
p.font.size = Pt(14)
p.font.color.rgb = TEXT_LIGHT
p.alignment = PP_ALIGN.CENTER

# Slide 8: REAL COMPETITOR PRICING DATA
slide = add_modern_content_slide(prs, "Market Competitive Analysis")

subtitle = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(8.8), Inches(0.3))
tf = subtitle.text_frame
p = tf.paragraphs[0]
p.text = "Real competitor pricing data: 1-800 Bollards, Global Industrial, Polector, U-Line"
p.font.size = Pt(11)
p.font.italic = True
p.font.color.rgb = GRAY

left = Inches(0.35)
top = Inches(1.55)
width = Inches(9.3)
height = Inches(5.4)

rows = 11
cols = 7

table = slide.shapes.add_table(rows, cols, left, top, width, height).table

table.columns[0].width = Inches(2.5)
table.columns[1].width = Inches(0.95)
table.columns[2].width = Inches(1.05)
table.columns[3].width = Inches(1.05)
table.columns[4].width = Inches(1.05)
table.columns[5].width = Inches(1.35)
table.columns[6].width = Inches(1.35)

headers = ["Product", "ZASP", "1-800", "Global", "Polector", "U-Line", "S4 New Price"]
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = PRIMARY
    para = cell.text_frame.paragraphs[0]
    para.font.size = Pt(10)
    para.font.bold = True
    para.font.color.rgb = TEXT_LIGHT
    para.alignment = PP_ALIGN.CENTER

for idx, (_, row) in enumerate(competitor.iterrows()):
    if pd.isna(row['Product Name']) or row['Product Name'] == '':
        continue

    table.cell(idx+1, 0).text = row['Product Name'][:30]
    table.cell(idx+1, 1).text = str(row['ZASP CPU']).replace('$', '').strip()
    table.cell(idx+1, 2).text = str(row['1800 CPU']).replace('$', '').strip()
    table.cell(idx+1, 3).text = str(row['GLOBAL']).replace('$', '').strip() if pd.notna(row['GLOBAL']) else '-'
    table.cell(idx+1, 4).text = str(row['POLECTOR']).replace('$', '').strip() if pd.notna(row['POLECTOR']) else '-'
    table.cell(idx+1, 5).text = str(row['ULINE']).replace('$', '').strip() if pd.notna(row['ULINE']) else '-'
    table.cell(idx+1, 6).text = str(row['S4 NEW MSRP']).replace('$', '').strip()

    for col in range(7):
        cell = table.cell(idx+1, col)
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(9)
        para.font.color.rgb = TEXT_DARK
        para.alignment = PP_ALIGN.CENTER if col > 0 else PP_ALIGN.LEFT

        if (idx+1) % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BG

        if col == 1:
            cell.fill.solid()
            cell.fill.fore_color.rgb = SUCCESS
            para.font.bold = True
            para.font.color.rgb = TEXT_LIGHT

        if col == 6:
            cell.fill.solid()
            cell.fill.fore_color.rgb = ACCENT
            para.font.bold = True
            para.font.color.rgb = TEXT_LIGHT

footnote = slide.shapes.add_textbox(Inches(0.35), Inches(7.05), Inches(9.3), Inches(0.35))
tf = footnote.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "ZASP pricing enables 50-80% margin improvement while remaining competitive with market leaders. Our new pricing maintains competitive positioning with significantly higher profitability."
p.font.size = Pt(9)
p.font.color.rgb = GRAY

# Slide 9: Next Steps
slide = add_modern_content_slide(prs, "Recommendation & Next Steps")

rec_banner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.4), Inches(8.6), Inches(1.2))
rec_banner.fill.solid()
rec_banner.fill.fore_color.rgb = SUCCESS
rec_banner.line.fill.background()
rec_banner.adjustments[0] = 0.05

rec_text = slide.shapes.add_textbox(Inches(0.9), Inches(1.6), Inches(8.2), Inches(0.8))
tf = rec_text.text_frame
p = tf.paragraphs[0]
p.text = "PROCEED WITH INITIAL ORDER"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = TEXT_LIGHT
p.alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = "198 units from ZASP for $33,599"
p.font.size = Pt(26)
p.font.bold = True
p.font.color.rgb = TEXT_LIGHT
p.alignment = PP_ALIGN.CENTER

benefits_title = slide.shapes.add_textbox(Inches(0.7), Inches(2.8), Inches(8.6), Inches(0.4))
tf = benefits_title.text_frame
p = tf.paragraphs[0]
p.text = "KEY BENEFITS"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = PRIMARY

benefits_content = slide.shapes.add_textbox(Inches(0.7), Inches(3.3), Inches(8.6), Inches(3))
tf = benefits_content.text_frame
tf.word_wrap = True

benefits = [
    "✓  62% cost reduction vs 1-800 Bollards distributor",
    "✓  Margin improvement from 23% to 72%",
    "✓  Additional $186K annual profit at current velocity",
    "✓  Investment payback in less than 2 months",
    "✓  3-6 months inventory coverage based on product velocity",
    "✓  Competitive market positioning with pricing flexibility"
]

for benefit in benefits:
    p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
    p.text = benefit
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_DARK
    p.space_after = Pt(10)

next_banner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(6.5), Inches(8.6), Inches(0.7))
next_banner.fill.solid()
next_banner.fill.fore_color.rgb = ACCENT
next_banner.line.fill.background()
next_banner.adjustments[0] = 0.05

next_text = slide.shapes.add_textbox(Inches(0.9), Inches(6.65), Inches(8.2), Inches(0.4))
tf = next_text.text_frame
p = tf.paragraphs[0]
p.text = "Next: Finalize ZASP agreement → Place initial order → Scale to container-level orders"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = TEXT_LIGHT
p.alignment = PP_ALIGN.CENTER

# Save presentation
output_path = r"C:\Users\blkw\OneDrive\Documents\Claude Code\Source 4 Industries\Bulk Bollards\Bulk_Bollard_Order_Proposal_FINAL_WITH_COMPETITORS.pptx"
prs.save(output_path)
print(f"Modern presentation with real competitor data created: {output_path}")
